import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import math
import logging
import re
import uuid

logger = logging.getLogger(__name__)

def is_total_row(row_value):
    """Enhanced TOTAL row detection function for all report types"""
    if pd.isna(row_value):
        return False
    
    first_col_value = str(row_value).upper().strip()
    
    # Comprehensive TOTAL detection patterns
    total_patterns = [
        'TOTAL', 'GRAND TOTAL', 'PART 1 TOTAL', 'PART 2 TOTAL',
        'OVERALL TOTAL', 'SUB TOTAL', 'SUBTOTAL', 'OVERALL',
        'SUM', 'SUMMARY', 'AGGREGATE', 'ALL EXECUTIVES',
        'COMBINED', 'GRAND', 'FINAL TOTAL', 'ACCLP',
        'ACCLLP', 'ACCLP TOTAL', 'ACCLLP TOTAL', 'ACCLP GRAND TOTAL'
    ]
    
    # Check exact matches
    if first_col_value in total_patterns:
        return True
    
    # Check if contains any total pattern
    for pattern in total_patterns:
        if pattern in first_col_value:
            return True
    
    # Check if ends with TOTAL
    if first_col_value.endswith('TOTAL'):
        return True
    
    return False

def format_column_header(header_text):
    """Format column headers with proper units and uppercase"""
    header_text = str(header_text).upper()
    
    # Enhanced label formatting for Product Growth
    if 'LY_QTY' in header_text:
        header_text = header_text.replace('LY_QTY', 'LAST-YEAR QTY/MT')
    elif 'LY_VALUE' in header_text:
        header_text = header_text.replace('LY_VALUE', 'LAST-YEAR VALUE/L')
    elif 'CY_QTY' in header_text:
        header_text = header_text.replace('CY_QTY', 'CURRENT-YEAR QTY/MT')
    elif 'CY_VALUE' in header_text:
        header_text = header_text.replace('CY_VALUE', 'CURRENT-YEAR VALUE/L')
    elif 'BUDGET_QTY' in header_text:
        header_text = header_text.replace('BUDGET_QTY', 'TARGET QTY/MT')
    elif 'BUDGET_VALUE' in header_text:
        header_text = header_text.replace('BUDGET_VALUE', 'TARGET VALUE/L')
    elif 'ACHIEVEMENT' in header_text:
        header_text = 'GROWTH %'
    elif '% ACHIEVED (SELECTED MONTH)' in header_text:
        header_text = 'FOR THE MONTH % ACHIEVED'
    elif 'BUDGET QTY' in header_text:
        header_text = header_text.replace('BUDGET QTY', 'TARGET QTY/MT')
    elif 'BUDGET VALUE' in header_text:
        header_text = header_text.replace('BUDGET VALUE', 'TARGET VALUE/L')
    elif header_text == 'BUDGET':
        header_text = 'TARGET'
    elif any(keyword in header_text for keyword in ['QTY', 'QUANTITY']):
        if '/MT' not in header_text:
            header_text = header_text.replace('QTY', 'QTY/MT').replace('QUANTITY', 'QUANTITY/MT')
    elif any(keyword in header_text for keyword in ['VALUE', 'AMOUNT', 'COLLECTION', 'TARGET', 'BILLED']):
        if 'ACHIEVEMENT' not in header_text and '%' not in header_text and 'GROWTH' not in header_text:
            if '/L' not in header_text:
                header_text += '/L'
    
    return header_text

def format_slide_title(title):
    """Remove quantity and value labels from slide titles and format properly"""
    title_lower = title.lower()
    
    # If it's a Product Growth title, preserve the company name and growth type
    if any(keyword in title_lower for keyword in ['quantity growth', 'value growth']):
        return title
    
    # Enhanced Product Growth title formatting
    title = title.replace(' LY:', ' Last-Year:').replace(' CY:', ' Current-Year:')
    title = title.replace('LY:', 'Last-Year:').replace('CY:', 'Current-Year:')
    
    # Remove common quantity and value indicators from titles
    patterns_to_remove = [
        r'\(Qty in Mt\)',
        r'\(Value in Lakhs\)',
        r'\(Qty in MT\)',
        r'\(Value in L\)',
    ]
    
    for pattern in patterns_to_remove:
        title = re.sub(pattern, '', title, flags=re.IGNORECASE)
    
    # Clean up extra spaces and dashes
    title = re.sub(r'\s+', ' ', title)
    title = re.sub(r'\s*-\s*$', '', title)
    title = title.strip()
    
    return title.upper()

def create_title_slide(prs, title):
    """Create enhanced title slide"""
    blank_slide_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(blank_slide_layout)
    
    # Company name
    company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    company_frame = company_name.text_frame
    company_frame.text = "ASIA CRYSTAL COMMODITY LLP"
    p = company_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    # Title
    formatted_title = format_slide_title(title)
    title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = formatted_title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    return title_slide

def add_slide_number(slide, slide_number):
    """Add slide number to the bottom right corner"""
    try:
        slide_number_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(6.8),
            Inches(1.5), Inches(0.5)
        )
        
        slide_number_frame = slide_number_box.text_frame
        slide_number_frame.text = str(slide_number)
        
        p = slide_number_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = "Times New Roman"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        
    except Exception as e:
        logger.error(f"Error adding slide number: {e}")

def create_thank_you_slide(prs):
    """Create thank you slide"""
    blank_slide_layout = prs.slide_layouts[6]
    thank_you_slide = prs.slides.add_slide(blank_slide_layout)
    
    thank_you_text = thank_you_slide.shapes.add_textbox(
        Inches(0.5), Inches(3.0),
        Inches(12.33), Inches(1.5)
    )
    thank_you_frame = thank_you_text.text_frame
    thank_you_frame.text = "THANK YOU"
    p = thank_you_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    return thank_you_slide

def determine_column_order(df, title):
    """Determine the correct column order based on report type"""
    actual_columns = list(df.columns)
    title_lower = title.lower()
    
    if any(phrase in title_lower for phrase in ["budget against billed", "budget vs billed", "budget v/s billed"]):
        has_qty_columns = any('qty' in col.lower() for col in actual_columns)
        has_value_columns = any('value' in col.lower() for col in actual_columns)
        
        if has_qty_columns and ("qty" in title_lower or "mt" in title_lower):
            return ['Executive', 'Budget Qty', 'Billed Qty', '%']
        elif has_value_columns and ("value" in title_lower or "lakhs" in title_lower):
            return ['Executive', 'Budget Value', 'Billed Value', '%']
        else:
            strict_order = ['Executive']
            budget_cols = [col for col in actual_columns if 'budget' in col.lower()]
            billed_cols = [col for col in actual_columns if 'billed' in col.lower()]
            percent_cols_detected = [col for col in actual_columns if '%' in col or 'percent' in col.lower()]
            
            strict_order.extend(budget_cols + billed_cols + percent_cols_detected)
            return strict_order
            
    elif any(phrase in title_lower for phrase in ["overall sales", "overall sale", "total sales"]):
        has_qty_columns = any('qty' in col.lower() for col in actual_columns)
        has_value_columns = any('value' in col.lower() for col in actual_columns)
        
        if has_qty_columns and ("qty" in title_lower or "mt" in title_lower):
            return ['Executive', 'Budget Qty', 'Billed Qty']
        elif has_value_columns:
            return ['Executive', 'Budget Value', 'Billed Value']
        else:
            strict_order = []
            if 'Executive' in actual_columns:
                strict_order.append('Executive')
            
            budget_cols = [col for col in actual_columns if 'budget' in col.lower() and col not in strict_order]
            strict_order.extend(budget_cols)
            
            billed_cols = [col for col in actual_columns if 'billed' in col.lower() and col not in strict_order]
            strict_order.extend(billed_cols)
            
            for col in actual_columns:
                if col not in strict_order:
                    strict_order.append(col)
            
            return strict_order
            
    elif any(phrase in title_lower for phrase in ["od target vs collection", "od collection"]):
        return [
            'Executive', 'Due Target', 'Collection Achieved',
            'Overall % Achieved', 'For the month Overdue',
            'For the month Collection', '% Achieved (Selected Month)'
        ]
        
    elif any(phrase in title_lower for phrase in ["product growth", "quantity growth"]):
        return ['PRODUCT GROUP', 'LY_QTY', 'BUDGET_QTY', 'CY_QTY', 'ACHIEVEMENT %']
        
    elif any(phrase in title_lower for phrase in ["value growth", "value in lakhs"]):
        return ['PRODUCT GROUP', 'LY_VALUE', 'BUDGET_VALUE', 'CY_VALUE', 'ACHIEVEMENT %']
        
    elif any(phrase in title_lower for phrase in ["customer", "billed customers"]):
        strict_order = []
        if 'S.No' in actual_columns: 
            strict_order.append('S.No')
        if 'Executive Name' in actual_columns: 
            strict_order.append('Executive Name')
        elif 'Executive' in actual_columns: 
            strict_order.append('Executive')
        for col in actual_columns:
            if col not in strict_order:
                strict_order.append(col)
        return strict_order
        
    else:
        return actual_columns.copy()

def add_enhanced_table_slide(prs, df, title, percent_cols=None, max_rows=14):
    """Add table slide with enhanced formatting and proper data handling"""
    if percent_cols is None:
        percent_cols = []

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title formatting
    formatted_title = format_slide_title(title)
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = formatted_title
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER

    # Determine column order
    ordered_columns = determine_column_order(df, title)
    
    # Ensure all expected columns exist
    final_columns = []
    for col in ordered_columns:
        if col in df.columns:
            final_columns.append(col)
    
    # Add any missing columns
    for col in df.columns:
        if col not in final_columns:
            final_columns.append(col)

    num_rows = len(df) + 1
    num_cols = len(final_columns)

    # Table sizing
    table_width = Inches(11.0)
    table_height = Inches(len(df) * 0.35 + 0.5)
    table_left = Inches(1.165)
    table_top = Inches(1.3)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        table_left, table_top,
        table_width, table_height
    ).table

    # Column width distribution
    if num_cols > 0:
        if num_cols == 2:
            table.columns[0].width = Inches(6.0)
            table.columns[1].width = Inches(5.0)
        elif 'Executive' in final_columns[0] or 'EXECUTIVE' in final_columns[0]:
            table.columns[0].width = Inches(3.5)
            remaining_width = 7.5
            if num_cols > 1:
                col_width = remaining_width / (num_cols - 1)
                for i in range(1, num_cols):
                    table.columns[i].width = Inches(col_width)
        else:
            col_width = 11.0 / num_cols
            for i in range(num_cols):
                table.columns[i].width = Inches(col_width)

    # Header row formatting
    for i, col_name in enumerate(final_columns):
        cell = table.cell(0, i)
        formatted_header = format_column_header(col_name)
        cell.text = formatted_header
        
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        
        para = cell.text_frame.paragraphs[0]
        para.font.name = "Calibri"
        para.font.size = Pt(15)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
        
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)

    # Data rows formatting
    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        first_col_value = row.iloc[0]
        is_total = is_total_row(first_col_value)
        
        for col_idx, col_name in enumerate(final_columns):
            cell = table.cell(row_idx, col_idx)
            value = row.get(col_name, "")
            
            # Format values based on column type
            if col_name in ['ACHIEVEMENT %', 'GROWTH %', '%', 'Overall % Achieved', '% Achieved (Selected Month)']:
                if isinstance(value, (int, float)) and not pd.isna(value):
                    cell.text = f"{value:.2f}%"
                else:
                    cell.text = str(value) if value is not None else ""
            elif col_idx in percent_cols and isinstance(value, (int, float)) and not pd.isna(value):
                cell.text = f"{value:.2f}%"
            elif isinstance(value, (int, float)) and not pd.isna(value):
                cell.text = str(int(round(value))) if abs(value - round(value)) < 0.001 else f"{value:.2f}"
            else:
                cell.text = str(value).upper() if value is not None and not pd.isna(value) else ""

            para = cell.text_frame.paragraphs[0]
            para.font.name = "Calibri"
            para.alignment = PP_ALIGN.CENTER
            
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            
            cell.fill.solid()

            # Enhanced row styling
            if is_total:
                cell.fill.fore_color.rgb = RGBColor(169, 169, 169)
                para.font.name = "Calibri"
                para.font.bold = True
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(255, 255, 255)
            else:
                para.font.bold = False
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(0, 0, 0)
                
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    return slide

def process_df_for_slides(prs, df, title_base, percent_cols=None, max_rows=14):
    """Process DataFrame and split into multiple slides if needed"""
    if percent_cols is None:
        percent_cols = []
    
    # Remove non-TOTAL "ACCLP" rows
    if not df.empty:
        first_col = df.columns[0]
        df = df[~((df[first_col] == "ACCLP") & (~df[first_col].apply(is_total_row)))].copy()
    
    # Separate data and total rows
    first_col = df.columns[0]
    total_mask = df[first_col].apply(is_total_row)
    data_rows = df[~total_mask].copy()
    total_row = df[total_mask].copy()
    
    num_data_rows = len(data_rows)
    
    if num_data_rows == 0:
        # Only TOTAL row(s)
        if not total_row.empty:
            add_enhanced_table_slide(prs, total_row, title_base, percent_cols=percent_cols, max_rows=15)
        return
    
    if num_data_rows <= max_rows:
        # Single slide
        if not total_row.empty:
            combined_df = pd.concat([data_rows, total_row], ignore_index=True)
        else:
            combined_df = data_rows.copy()
        add_enhanced_table_slide(prs, combined_df, title_base, percent_cols=percent_cols, max_rows=15)
        return
    
    # Split into multiple slides
    num_parts = math.ceil(num_data_rows / max_rows)
    
    for i in range(num_parts):
        start_idx = i * max_rows
        end_idx = min((i + 1) * max_rows, num_data_rows)
        part = data_rows.iloc[start_idx:end_idx].copy()
        
        # Add total row to last part
        is_last_part = (i == num_parts - 1)
        
        if is_last_part and not total_row.empty:
            combined_part = pd.concat([part, total_row], ignore_index=True)
        else:
            combined_part = part.copy()
        
        # Round numeric columns
        numeric_cols = combined_part.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            combined_part[col] = combined_part[col].apply(lambda x: round(float(x), 2) if pd.notna(x) else x)
        
        add_enhanced_table_slide(prs, combined_part, title_base, percent_cols=percent_cols, max_rows=15)

def generate_consolidated_ppt_enhanced(payload):
    """Enhanced consolidated PPT generation with proper formatting and data handling"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title = payload.get("title", "Consolidated Report")
    dfs_info = payload.get("dfs_with_titles", [])

    # Create title slide
    create_title_slide(prs, title)

    logger.info(f"Creating consolidated PPT with {len(dfs_info)} reports")

    # Process each report
    for df_info in dfs_info:
        df_data = df_info.get('data', [])
        slide_title = df_info['title']
        percent_cols = df_info.get('percent_cols', [])
        columns = df_info.get('columns', [])
        
        # Convert to DataFrame
        if isinstance(df_data, list):
            df = pd.DataFrame(df_data)
        else:
            df = df_data.copy()

        # Apply column order if specified
        if columns:
            df = df[[col for col in columns if col in df.columns]]
        
        if df.empty:
            logger.warning(f"Skipping empty report: {slide_title}")
            continue
        
        # Clean DataFrame
        df = df.dropna(how='all')
        
        # Fill NaN values
        for col in df.columns:
            if df[col].dtype == 'object':
                df[col] = df[col].fillna('')
            else:
                df[col] = df[col].fillna(0)
        
        # Round numeric columns
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            df[col] = df[col].round(2)
        
        logger.info(f"Processing consolidated report: {slide_title}")
        
        # Process with enhanced splitting
        process_df_for_slides(prs, df, slide_title, percent_cols=percent_cols, max_rows=14)

    # Create thank you slide
    create_thank_you_slide(prs)
    
    # Add slide numbers (skip title slide)
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if i == 0:  # Skip title slide
            continue
        else:
            add_slide_number(slide, i)

    # Save presentation
    out_path = f"static/Enhanced_Consolidated_Report_{uuid.uuid4().hex[:6]}.pptx"
    prs.save(out_path)
    
    logger.info(f"Enhanced consolidated PPT created with {total_slides} slides")
    return out_path