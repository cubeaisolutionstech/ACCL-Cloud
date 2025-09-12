import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO
import math
import logging
import base64
import re
from datetime import datetime
import os
import uuid
import json

logger = logging.getLogger(__name__)

# Standardized color scheme
class PPTColors:
    """Centralized color configuration for PPT generation"""
    
    # Header colors
    HEADER_BG = RGBColor(0, 112, 192)  # Blue header background
    HEADER_TEXT = RGBColor(255, 255, 255)  # White header text
    
    # Total row colors - GRAY SCHEME
    TOTAL_ROW_BG = RGBColor(211, 211, 211)  # Gray background (matching Streamlit)
    TOTAL_ROW_TEXT = RGBColor(0, 0, 0)  # Black text
    
    # Regular data row colors
    ALTERNATING_ROW_1 = RGBColor(240, 248, 255)  # Light blue (matching Streamlit)
    ALTERNATING_ROW_2 = RGBColor(255, 255, 255)  # White
    REGULAR_TEXT = RGBColor(0, 0, 0)             # Black text

def is_total_row(row_value):
    """Enhanced TOTAL row detection function matching Streamlit logic"""
    if pd.isna(row_value):
        return False
    
    first_col_value = str(row_value).upper().strip()
    
    # Comprehensive TOTAL detection patterns from Streamlit
    total_patterns = [
        'TOTAL', 'GRAND TOTAL', 'PART 1 TOTAL', 'PART 2 TOTAL',
        'OVERALL TOTAL', 'SUB TOTAL', 'SUBTOTAL', 'OVERALL',
        'SUM', 'SUMMARY', 'AGGREGATE', 'ALL EXECUTIVES',
        'COMBINED', 'GRAND', 'FINAL TOTAL', 'ACCLP',
        'ACCLLP', 'ACCLP TOTAL', 'ACCLLP TOTAL', 'ACCLP GRAND TOTAL'
    ]
    
    # Check exact matches
    if first_col_value in total_patterns:
        return True
    
    # Check if contains any total pattern
    for pattern in total_patterns:
        if pattern in first_col_value:
            return True
    
    # Check if ends with TOTAL
    if first_col_value.endswith('TOTAL'):
        return True
    
    return False

def add_page_number(slide, page_num):
    """Add page number to the bottom right of a slide - Streamlit style"""
    try:
        # Add page number text box at bottom right
        page_number_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3)
        )
        page_frame = page_number_box.text_frame
        page_frame.text = str(page_num)
        
        # Format page number - Streamlit style
        p = page_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = False
        p.font.color.rgb = RGBColor(128, 128, 128)  # Gray color
    except Exception as e:
        logger.error(f"Error adding page number: {e}")

def create_title_slide(prs, title, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create title slide matching Streamlit format"""
    try:
        blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
        title_slide = prs.slides.add_slide(blank_slide_layout)
        
        # Company name - matching Streamlit style
        company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        company_frame = company_name.text_frame
        company_frame.text = "ASIA CRYSTAL COMMODITY LLP"
        p = company_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        
        # Logo handling - Streamlit style
        if logo_file is not None:
            try:
                if isinstance(logo_file, str):
                    # Base64 encoded logo
                    logo_data = base64.b64decode(logo_file)
                    logo_buffer = BytesIO(logo_data)
                else:
                    # File object or BytesIO
                    logo_buffer = BytesIO(logo_file.read())
                    logo_file.seek(0)  # Reset for reuse
                
                title_slide.shapes.add_picture(
                    logo_buffer, Inches(5.665), Inches(1.5), 
                    width=Inches(2), height=Inches(2)
                )
            except Exception as e:
                logger.error(f"Error adding logo to slide: {str(e)}")
        
        # Title - matching Streamlit format
        title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title.upper()  # Uppercase like Streamlit
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        # Executive Information Section (if provided)
        if executive_name or date_str or branch_name:
            info_y_position = 6.0
            
            if executive_name:
                name_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(info_y_position), Inches(12.33), Inches(0.5))
                name_frame = name_box.text_frame
                name_frame.text = f"NAME: {executive_name.upper()}"
                p = name_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.name = "Times New Roman"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(128, 0, 128)
                info_y_position += 0.4
            
            if date_str:
                date_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(info_y_position), Inches(12.33), Inches(0.5))
                date_frame = date_box.text_frame
                date_frame.text = f"DATE: {date_str}"
                p = date_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.name = "Times New Roman"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(128, 0, 128)
        
        return title_slide
        
    except Exception as e:
        logger.error(f"Error creating title slide: {e}")
        raise

def create_thank_you_slide(prs, logo_file=None):
    """Create thank you slide matching Streamlit format"""
    try:
        blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
        thank_you_slide = prs.slides.add_slide(blank_slide_layout)
        
        # Main "Thank You" text - matching Streamlit positioning
        thank_you_box = thank_you_slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.33), Inches(1.5))
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.text = "THANK YOU"
        p = thank_you_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        return thank_you_slide
        
    except Exception as e:
        logger.error(f"Error creating thank you slide: {e}")
        raise

def add_table_slide(prs, df, title, percent_cols=None, page_num=None):
    """Add table slide matching Streamlit format exactly"""
    if df is None or df.empty:
        return None
    
    if percent_cols is None:
        percent_cols = []
    
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title with standardized font - matching Streamlit
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.75))
    title_frame = title_shape.text_frame
    title_frame.text = title.upper()  # Ensure title is uppercase
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    # Table setup - matching Streamlit dimensions
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5)).table
    
    # Header row with uppercase text - matching Streamlit
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name).upper()  # Convert to uppercase
        cell.text_frame.paragraphs[0].font.size = Pt(15)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PPTColors.HEADER_BG
        cell.text_frame.paragraphs[0].font.color.rgb = PPTColors.HEADER_TEXT
    
    # Data rows - matching Streamlit logic exactly
    for row_idx in range(df.shape[0]):
        # Check if this is a total row - using Streamlit logic
        is_total_row_check = is_total_row(df.iloc[row_idx, 0])
        
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            value = df.iloc[row_idx, col_idx]
            
            # Format cells with consistent decimal places - matching Streamlit
            if percent_cols and col_idx in percent_cols:
                cell.text = f"{value}%"
            else:
                # For numeric columns (not first column), ensure 2-decimal formatting
                if col_idx > 0 and isinstance(value, (int, float)) and not pd.isna(value):
                    if "%" not in str(df.columns[col_idx]):
                        cell.text = f"{float(value):.2f}"
                    else:
                        cell.text = str(value)
                elif col_idx == 0:
                    cell.text = str(value).upper()
                else:
                    cell.text = str(value)
            
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # Apply Streamlit-style formatting
            if is_total_row_check:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = PPTColors.TOTAL_ROW_BG
                cell.text_frame.paragraphs[0].font.color.rgb = PPTColors.TOTAL_ROW_TEXT
            else:
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PPTColors.ALTERNATING_ROW_1
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PPTColors.ALTERNATING_ROW_2
                cell.text_frame.paragraphs[0].font.color.rgb = PPTColors.REGULAR_TEXT
    
    # Add page number if provided
    if page_num is not None:
        add_page_number(slide, page_num)
    
    return slide

def format_budget_dataframes(budget_vs_billed_qty_df, budget_vs_billed_value_df, 
                            overall_sales_qty_df, overall_sales_value_df):
    """Format DataFrames with proper column names and order for budget reports"""

    # Define mappings (you can adjust based on your raw data)
    qty_mapping = {
        'AREA': 'Area',
        'Budget Qty': 'Budget Qty/Mt',
        'Billed Qty': 'Billed Qty/Mt',
        '%': '%'
    }
    value_mapping = {
        'AREA': 'Area',
        'Budget Value': 'Budget Value/L',
        'Billed Value': 'Billed Value/L',
        '%': '%'
    }

    def clean_df(df, mapping, expected_cols):
        if df.empty:
            return df
        # Standardize column names (strip spaces, make upper for match)
        df.columns = [str(c).strip().upper() for c in df.columns]
        # Map to friendly names
        df = df.rename(columns={k.upper(): v for k, v in mapping.items()})
        # Reorder only if all exist
        available_cols = [c for c in expected_cols if c in df.columns]
        return df[available_cols]

    budget_vs_billed_qty_df = clean_df(budget_vs_billed_qty_df, qty_mapping, ['Area', 'Budget Qty/Mt', 'Billed Qty/Mt', '%'])
    budget_vs_billed_value_df = clean_df(budget_vs_billed_value_df, value_mapping, ['Area', 'Budget Value/L', 'Billed Value/L', '%'])
    overall_sales_qty_df = clean_df(overall_sales_qty_df, qty_mapping, ['Area', 'Budget Qty/Mt', 'Billed Qty/Mt'])
    overall_sales_value_df = clean_df(overall_sales_value_df, value_mapping, ['Area', 'Budget Value/L', 'Billed Value/L'])

    return budget_vs_billed_qty_df, budget_vs_billed_value_df, overall_sales_qty_df, overall_sales_value_df

def create_budget_ppt(budget_vs_billed_qty_df, budget_vs_billed_value_df, overall_sales_qty_df, overall_sales_value_df, month_title=None, logo_file=None):
    """Create PPT presentation with properly formatted columns and titles"""
    try:
        # Note: DataFrames should already be formatted when passed to this function
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide (no page number) - matching Streamlit
        title = f"Monthly Review Meeting – {month_title.upper()}" if month_title else "Budget vs Billed Analysis"
        create_title_slide(prs, title, logo_file)
        
        page_num = 1  # Start page numbering from 1 for content slides
        
        # Add slides with consistent titles (without units in title since they're in columns now)
        def add_table_slide_improved(title_text, df, page_number, percent_cols=None):
            if percent_cols is None:
                percent_cols = []
            slide = add_table_slide(prs, df, title_text, percent_cols=percent_cols, page_num=page_number)
            return slide
        
        # Budget vs Billed slides - titles without units since columns have units
        add_table_slide_improved(f"BUDGET VS BILLED - QUANTITY - {month_title.upper()}", budget_vs_billed_qty_df, page_num, percent_cols=[3])
        page_num += 1
        add_table_slide_improved(f"BUDGET VS BILLED - VALUE - {month_title.upper()}", budget_vs_billed_value_df, page_num, percent_cols=[3])
        page_num += 1
        add_table_slide_improved(f"OVERALL SALES - QUANTITY - {month_title.upper()}", overall_sales_qty_df, page_num)
        page_num += 1
        add_table_slide_improved(f"OVERALL SALES - VALUE - {month_title.upper()}", overall_sales_value_df, page_num)
        
        # Add thank you slide (no page number)
        create_thank_you_slide(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating Budget PPT: {e}")
        return None
def format_od_target_vs_collection_dataframes(branch_df, regional_df):
    """Format OD Target vs Collection DataFrames with proper column names matching Streamlit exactly"""
    
    # Define proper column orders and names with units matching Streamlit
    od_branch_columns = ['BRANCH', 'DUE TARGET/L', 'COLLECTION ACHIEVED/L', 'OVERALL % ACHIEVED', 
                        'FOR THE MONTH OVERDUE/L', 'FOR THE MONTH COLLECTION/L', 'FOR THE MONTH % ACHIEVED']
    od_regional_columns = ['REGION', 'DUE TARGET/L', 'COLLECTION ACHIEVED/L', 'OVERALL % ACHIEVED',
                          'FOR THE MONTH OVERDUE/L', 'FOR THE MONTH COLLECTION/L', 'FOR THE MONTH % ACHIEVED']
    
    # Format Branch DataFrame
    if not branch_df.empty:
        current_cols = branch_df.columns.tolist()
        
        if len(current_cols) >= 7:
            expected_mapping = {
                current_cols[0]: 'BRANCH',
                current_cols[1]: 'DUE TARGET/L',
                current_cols[2]: 'COLLECTION ACHIEVED/L', 
                current_cols[3]: 'OVERALL % ACHIEVED',
                current_cols[4]: 'FOR THE MONTH OVERDUE/L',
                current_cols[5]: 'FOR THE MONTH COLLECTION/L',
                current_cols[6]: 'FOR THE MONTH % ACHIEVED'
            }
            
            branch_df = branch_df.rename(columns=expected_mapping)
            available_cols = [col for col in od_branch_columns if col in branch_df.columns]
            if available_cols:
                branch_df = branch_df[available_cols]
    
    # Format Regional DataFrame
    if regional_df is not None and not regional_df.empty:
        current_cols = regional_df.columns.tolist()
        
        # CRITICAL FIX: Handle multiple possible percentage column names
        if len(current_cols) >= 7:
            # Create flexible mapping that handles different percentage column names
            expected_mapping = {}
            
            for i, col in enumerate(current_cols):
                if i == 0:
                    expected_mapping[col] = 'REGION'
                elif i == 1:
                    expected_mapping[col] = 'DUE TARGET/L'
                elif i == 2:
                    expected_mapping[col] = 'COLLECTION ACHIEVED/L'
                elif i == 3:
                    expected_mapping[col] = 'OVERALL % ACHIEVED'
                elif i == 4:
                    expected_mapping[col] = 'FOR THE MONTH OVERDUE/L'
                elif i == 5:
                    expected_mapping[col] = 'FOR THE MONTH COLLECTION/L'
                elif i == 6:
                    # Handle both possible percentage column names
                    expected_mapping[col] = 'FOR THE MONTH % ACHIEVED'
            
            # Also handle specific known variations
            known_variations = {
                '% Achieved (Selected Month)': 'FOR THE MONTH % ACHIEVED',
                'For the month % Achieved': 'FOR THE MONTH % ACHIEVED'
            }
            expected_mapping.update(known_variations)
            
            regional_df = regional_df.rename(columns=expected_mapping)
            
            # Ensure proper column order
            available_cols = [col for col in od_regional_columns if col in regional_df.columns]
            if available_cols:
                regional_df = regional_df[available_cols]
    
    return branch_df, regional_df
def create_od_ppt_updated(branch_df, regional_df, title, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create OD PPT with properly formatted columns matching Streamlit exactly"""
    try:
        # Format DataFrames with proper column names and order matching Streamlit
        branch_df, regional_df = format_od_target_vs_collection_dataframes(branch_df, regional_df)
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide (no page number)
        create_title_slide(prs, title.upper(), logo_file, executive_name, date_str, branch_name)
        
        page_num = 1
        
        # Branch-wise performance slide
        if not branch_df.empty:
            # Clean title - remove units since they're now in column headers
            clean_title = title.replace(' - Value in Lakhs', '').replace(' - Qty in Mt', '')
            slide = add_table_slide(prs, branch_df, f"OD TARGET VS COLLECTION - BRANCH - {clean_title.split(' - ')[-1].upper()}", 
                                  percent_cols=[3, 6], page_num=page_num)
            if slide:
                page_num += 1
        
        # Regional summary slide (if exists)  
        if regional_df is not None and not regional_df.empty:
            clean_title = title.replace(' - Value in Lakhs', '').replace(' - Qty in Mt', '')
            slide = add_table_slide(prs, regional_df, f"OD TARGET VS COLLECTION - REGIONAL - {clean_title.split(' - ')[-1].upper()}", 
                                  percent_cols=[3, 6], page_num=page_num)
        
        # Add thank you slide (no page number)
        create_thank_you_slide(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating OD PPT: {e}")
        return None

def format_product_growth_dataframes(group_results):
    """Format Product Growth DataFrames with proper column names and order matching individual reports - FIXED VERSION"""
    
    # Define proper column orders for Product Growth reports
    qty_columns = ['PRODUCT NAME', 'LAST YEAR QTY/MT', 'BUDGET QTY/MT', 'CURRENT YEAR QTY/MT', 'GROWTH %']
    value_columns = ['PRODUCT NAME', 'LAST YEAR VALUE/L', 'BUDGET VALUE/L', 'CURRENT YEAR VALUE/L', 'GROWTH %']
    
    formatted_results = {}
    
    for group, data in group_results.items():
        # Format quantity DataFrame
        qty_df = data['qty_df'].copy() if isinstance(data['qty_df'], pd.DataFrame) else pd.DataFrame(data['qty_df'])
        value_df = data['value_df'].copy() if isinstance(data['value_df'], pd.DataFrame) else pd.DataFrame(data['value_df'])
        
        # Debug: Print original column names
        logger.info(f"Original qty_df columns for {group}: {qty_df.columns.tolist()}")
        logger.info(f"Original value_df columns for {group}: {value_df.columns.tolist()}")
        
        # Format quantity DataFrame with proper column mapping
        if not qty_df.empty:
            # Create a standardized mapping based on expected backend column names
            qty_column_mapping = {
                'PRODUCT NAME': 'PRODUCT NAME',
                'LAST_YEAR_QTY/MT': 'LAST YEAR QTY/MT',
                'LAST YEAR QTY/MT': 'LAST YEAR QTY/MT',
                'BUDGET_QTY/MT': 'BUDGET QTY/MT', 
                'BUDGET QTY/MT': 'BUDGET QTY/MT',
                'CURRENT_YEAR_QTY/MT': 'CURRENT YEAR QTY/MT',
                'CURRENT YEAR QTY/MT': 'CURRENT YEAR QTY/MT',
                'GROWTH %': 'GROWTH %'
            }
            
            # Apply column mapping
            qty_df = qty_df.rename(columns=qty_column_mapping)
            
            # Ensure we have all required columns, fill missing ones with 0
            for col in qty_columns:
                if col not in qty_df.columns:
                    logger.warning(f"Missing column {col} in qty_df for {group}, adding with 0 values")
                    qty_df[col] = 0 if col != 'PRODUCT NAME' else 'Unknown'
            
            # Reorder columns to match expected order
            qty_df = qty_df[qty_columns]
            
            logger.info(f"Formatted qty_df columns for {group}: {qty_df.columns.tolist()}")
        
        # Format value DataFrame with proper column mapping
        if not value_df.empty:
            # Create a standardized mapping based on expected backend column names
            value_column_mapping = {
                'PRODUCT NAME': 'PRODUCT NAME',
                'LAST_YEAR_VALUE/L': 'LAST YEAR VALUE/L',
                'LAST YEAR VALUE/L': 'LAST YEAR VALUE/L',
                'BUDGET_VALUE/L': 'BUDGET VALUE/L',
                'BUDGET VALUE/L': 'BUDGET VALUE/L',
                'CURRENT_YEAR_VALUE/L': 'CURRENT YEAR VALUE/L',
                'CURRENT YEAR VALUE/L': 'CURRENT YEAR VALUE/L',
                'GROWTH %': 'GROWTH %'
            }
            
            # Apply column mapping
            value_df = value_df.rename(columns=value_column_mapping)
            
            # Ensure we have all required columns, fill missing ones with 0
            for col in value_columns:
                if col not in value_df.columns:
                    logger.warning(f"Missing column {col} in value_df for {group}, adding with 0 values")
                    value_df[col] = 0 if col != 'PRODUCT NAME' else 'Unknown'
            
            # Reorder columns to match expected order
            value_df = value_df[value_columns]
            
            logger.info(f"Formatted value_df columns for {group}: {value_df.columns.tolist()}")
        
        formatted_results[group] = {
            'qty_df': qty_df,
            'value_df': value_df
        }
    
    return formatted_results

def create_product_growth_ppt(group_results, month_title, logo_file=None, ly_month=None, cy_month=None, executive_name=None, date_str=None, branch_name=None):
    """Create Product Growth PPT with properly formatted columns and clean titles"""
    try:
        # Format all DataFrames with proper column names and order
        formatted_results = format_product_growth_dataframes(group_results)
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide (no page number) - clean title without LY/CY details
        clean_title = f"PRODUCT GROWTH REPORT"
        create_title_slide(prs, clean_title, logo_file, executive_name, date_str, branch_name)
        
        page_num = 1
        
        # Add data slides for each company group
        for group, data in formatted_results.items():
            qty_df = data['qty_df']
            value_df = data['value_df']
            
            # Quantity slide - clean title without extra month details
            slide = add_table_slide(prs, qty_df, f"{group.upper()} - QUANTITY GROWTH", percent_cols=[4], page_num=page_num)
            if slide:
                page_num += 1
            
            # Value slide - clean title without extra month details
            slide = add_table_slide(prs, value_df, f"{group.upper()} - VALUE GROWTH", percent_cols=[4], page_num=page_num)
            if slide:
                page_num += 1
        
        # Add thank you slide (no page number)
        create_thank_you_slide(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating Product Growth PPT: {e}")
        return None

def create_nbc_individual_ppt(customer_df, customer_title, sorted_months, financial_year, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create NBC PPT with proper column handling"""
    try:
        from utils.ppt_generator import (
            create_title_slide, create_thank_you_slide, add_table_slide, 
            is_total_row, PPTColors
        )
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Remove S.No variations if they exist
        sno_variations = ['S.No', 'SNo', 'S No', 'sno', 's.no', 'Serial No', 'Sr. No']
        for sno_col in sno_variations:
            if sno_col in customer_df.columns:
                customer_df = customer_df.drop(sno_col, axis=1)
        
        # Ensure Branch column is correctly named
        if 'Mapped_Branch' in customer_df.columns:
            customer_df = customer_df.rename(columns={'Mapped_Branch': 'Branch'})
        
        # Validate we have Branch column
        if 'Branch' not in customer_df.columns:
            raise ValueError("Branch column not found in DataFrame")
        
        # Check for null values in Branch column
        null_branches = customer_df['Branch'].isna().sum()
        if null_branches > 0:
            logger.warning(f"Found {null_branches} null values in Branch column")
            customer_df['Branch'] = customer_df['Branch'].fillna('UNKNOWN')
        
        # Extract month for title
        month_for_title = ""
        if " - " in customer_title:
            month_for_title = customer_title.split(' - ', 1)[1].strip()
        elif sorted_months and len(sorted_months) > 0:
            month_for_title = sorted_months[-1]
        
        title = f"NUMBER OF BILLED CUSTOMERS - {month_for_title}" if month_for_title else "NUMBER OF BILLED CUSTOMERS"
        
        # Create slides
        create_title_slide(prs, title, logo_file, executive_name, date_str, branch_name)
        
        if not customer_df.empty:
            add_table_slide(prs, customer_df, title, page_num=1)
        
        create_thank_you_slide(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        logger.info(f"NBC PPT created successfully - Title: {title}")
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating NBC PPT: {e}")
        return None
def create_od_individual_ppt(od_target_df, od_title, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create OD Target PPT matching Streamlit format exactly - FIXED column handling"""
    try:
        # Make a copy to avoid modifying original DataFrame
        df = od_target_df.copy()
        
        # FIXED: Ensure column is named "TARGET/L" for display (matching Streamlit exactly)
        if 'TARGET' in df.columns and 'TARGET/L' not in df.columns:
            df = df.rename(columns={'TARGET': 'TARGET/L'})
        elif 'Area' in df.columns and len(df.columns) == 2:
            # If we have Area and one other column, rename it to TARGET/L
            other_col = [col for col in df.columns if col != 'Area'][0]
            df = df.rename(columns={other_col: 'TARGET/L'})
        
        # Ensure we have the right column structure
        if 'Area' not in df.columns:
            if 'AREA' in df.columns:
                df = df.rename(columns={'AREA': 'Area'})
        
        # Create PPT
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide (no page number)
        create_title_slide(prs, "OD TARGET REPORT", logo_file, executive_name, date_str, branch_name)
        
        # Add OD Target slide with page number
        slide = add_table_slide(prs, df, od_title.upper(), page_num=1)
        
        # Add thank you slide (no page number)
        create_thank_you_slide(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating OD Target PPT: {e}")
        return None

def create_consolidated_ppt(all_dfs_with_titles, logo_file=None, title="ACCLLP Consolidated Report", executive_name=None, date_str=None, branch_name=None):
    """Create consolidated PPT matching Streamlit format exactly with proper budget formatting"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Auto-load logo from static folder if not provided
        if logo_file is None:
            static_logo_path = os.path.join("static", "logo.jpeg")
            if os.path.exists(static_logo_path):
                try:
                    with open(static_logo_path, "rb") as f:
                        logo_file = BytesIO(f.read())
                        logo_file.seek(0)
                    logger.info(f"Logo loaded from {static_logo_path}")
                except Exception as e:
                    logger.warning(f"Could not read logo from static: {e}")
                    logo_file = None
        
        # Create title slide (no page number)
        create_title_slide(prs, title, logo_file, executive_name, date_str, branch_name)
        
        page_num = 1  # Start page numbering from 1 for content slides
        
        # Add content slides with page numbers - matching Streamlit
        for df_info in all_dfs_with_titles:
            if df_info and 'df' in df_info and 'title' in df_info:
                df_data = df_info['df']
                slide_title = df_info['title']
                percent_cols = df_info.get('percent_cols', [])
                report_type = df_info.get('report_type', '')  # Get report type if available
                
                # Convert to DataFrame if it's a list of dicts
                if isinstance(df_data, list):
                    df = pd.DataFrame(df_data)
                else:
                    df = df_data.copy()
                
                if df.empty:
                    logger.warning(f"Skipping empty report: {slide_title}")
                    continue
                
                # **CRITICAL FIX**: Apply proper formatting based on report type
                if 'BUDGET' in slide_title.upper() or 'budget' in report_type.lower():
                    # This is a budget report - ensure proper budget column formatting
                    logger.info(f"Processing budget report: {slide_title}")
                    
                    # Apply budget-specific column formatting
                    if 'QUANTITY' in slide_title.upper():
                        # Budget Quantity formatting
                        budget_qty_mapping = {
                            'AREA': 'Area',
                            'Area': 'Area',
                            'Budget Qty': 'Budget Qty/Mt',
                            'Budget_Qty': 'Budget Qty/Mt',
                            'BUDGET_QTY': 'Budget Qty/Mt',
                            'Billed Qty': 'Billed Qty/Mt',
                            'Billed_Qty': 'Billed Qty/Mt',
                            'BILLED_QTY': 'Billed Qty/Mt',
                            '%': '%',
                            'Percentage': '%',
                            'PCT': '%'
                        }
                        df = df.rename(columns=budget_qty_mapping)
                        expected_cols = ['Area', 'Budget Qty/Mt', 'Billed Qty/Mt', '%']
                        
                    elif 'VALUE' in slide_title.upper():
                        # Budget Value formatting
                        budget_value_mapping = {
                            'AREA': 'Area',
                            'Area': 'Area',
                            'Budget Value': 'Budget Value/L',
                            'Budget_Value': 'Budget Value/L',
                            'BUDGET_VALUE': 'Budget Value/L',
                            'Billed Value': 'Billed Value/L',
                            'Billed_Value': 'Billed Value/L',
                            'BILLED_VALUE': 'Billed Value/L',
                            '%': '%',
                            'Percentage': '%',
                            'PCT': '%'
                        }
                        df = df.rename(columns=budget_value_mapping)
                        expected_cols = ['Area', 'Budget Value/L', 'Billed Value/L', '%']
                    
                    # Reorder columns if they exist
                    available_cols = [col for col in expected_cols if col in df.columns]
                    if available_cols:
                        df = df[available_cols]
                        logger.info(f"Budget report formatted with columns: {df.columns.tolist()}")
                
                elif 'OD TARGET' in slide_title.upper() or 'od' in report_type.lower():
                    # OD Target specific formatting
                    logger.info(f"Processing OD report: {slide_title}")
                    
                    # Apply OD-specific column formatting
                    od_mapping = {
                        'BRANCH': 'BRANCH',
                        'Branch': 'BRANCH',
                        'DUE TARGET': 'DUE TARGET/L',
                        'DUE_TARGET': 'DUE TARGET/L',
                        'TARGET': 'DUE TARGET/L',
                        'COLLECTION ACHIEVED': 'COLLECTION ACHIEVED/L',
                        'COLLECTION_ACHIEVED': 'COLLECTION ACHIEVED/L',
                        'OVERALL % ACHIEVED': 'OVERALL % ACHIEVED',
                        'FOR THE MONTH OVERDUE': 'FOR THE MONTH OVERDUE/L',
                        'FOR_THE_MONTH_OVERDUE': 'FOR THE MONTH OVERDUE/L',
                        'FOR THE MONTH COLLECTION': 'FOR THE MONTH COLLECTION/L',
                        'FOR_THE_MONTH_COLLECTION': 'FOR THE MONTH COLLECTION/L',
                        'FOR THE MONTH % ACHIEVED': 'FOR THE MONTH % ACHIEVED'
                    }
                    df = df.rename(columns=od_mapping)
                
                elif 'PRODUCT GROWTH' in slide_title.upper() or 'product' in report_type.lower():
                    # Product Growth specific formatting
                    logger.info(f"Processing Product Growth report: {slide_title}")
                    
                    if 'QUANTITY' in slide_title.upper():
                        pg_qty_mapping = {
                            'PRODUCT NAME': 'PRODUCT NAME',
                            'Product_Name': 'PRODUCT NAME',
                            'LAST YEAR QTY': 'LAST YEAR QTY/MT',
                            'LAST_YEAR_QTY': 'LAST YEAR QTY/MT',
                            'BUDGET QTY': 'BUDGET QTY/MT',
                            'BUDGET_QTY': 'BUDGET QTY/MT',
                            'CURRENT YEAR QTY': 'CURRENT YEAR QTY/MT',
                            'CURRENT_YEAR_QTY': 'CURRENT YEAR QTY/MT',
                            'GROWTH %': 'GROWTH %'
                        }
                        df = df.rename(columns=pg_qty_mapping)
                    elif 'VALUE' in slide_title.upper():
                        pg_value_mapping = {
                            'PRODUCT NAME': 'PRODUCT NAME',
                            'Product_Name': 'PRODUCT NAME',
                            'LAST YEAR VALUE': 'LAST YEAR VALUE/L',
                            'LAST_YEAR_VALUE': 'LAST YEAR VALUE/L',
                            'BUDGET VALUE': 'BUDGET VALUE/L',
                            'BUDGET_VALUE': 'BUDGET VALUE/L',
                            'CURRENT YEAR VALUE': 'CURRENT YEAR VALUE/L',
                            'CURRENT_YEAR_VALUE': 'CURRENT YEAR VALUE/L',
                            'GROWTH %': 'GROWTH %'
                        }
                        df = df.rename(columns=pg_value_mapping)
                
                # Add the slide with proper formatting
                slide = add_table_slide(prs, df, slide_title, percent_cols=percent_cols, page_num=page_num)
                if slide:
                    page_num += 1
                    logger.info(f"Added slide #{page_num-1}: {slide_title}")
        
        # Add thank you slide (no page number)
        create_thank_you_slide(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        logger.info(f"Consolidated PPT created with {len(prs.slides)} total slides")
        return ppt_buffer
    
    except Exception as e:
        logger.error(f"Error creating consolidated PPT: {e}")
        return None

# Export all functions
__all__ = [
    'PPTColors',
    'is_total_row',
    'add_page_number',
    'create_title_slide',
    'create_thank_you_slide',
    'add_table_slide',
    'create_budget_ppt',
    'create_od_ppt_updated',
    'create_product_growth_ppt',
    'create_nbc_individual_ppt',
    'create_od_individual_ppt',
    'create_consolidated_ppt'
]
