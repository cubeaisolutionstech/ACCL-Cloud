import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import math
import logging
import base64
import re
from datetime import datetime

logger = logging.getLogger(__name__)

def calculate_slide_constraints():
    """Calculate exact slide constraints for table placement"""
    slide_width = 13.33  # inches
    slide_height = 7.5   # inches
    
    # Reserved areas
    title_area_height = 1.1  # Title + margins
    slide_number_area = 0.3  # Bottom area for slide numbers
    side_margins = 0.8      # REDUCED from 1.0 to 0.8 for wider tables
    
    # Available area for table
    available_width = slide_width - side_margins
    available_height = slide_height - title_area_height - slide_number_area
    
    return {
        'slide_width': slide_width,
        'slide_height': slide_height,
        'available_width': available_width,
        'available_height': available_height,
        'table_left': 0.4,  # REDUCED from 0.5 to 0.4
        'table_top': 1.1,   # Below title
        'max_table_width': available_width,
        'max_table_height': available_height
    }

def calculate_max_fitting_rows():
    """Calculate maximum rows that can fit with minimum readable size - FIXED TO 15 ROWS"""
    constraints = calculate_slide_constraints()
    
    # FIXED: Always return 14 data rows max per slide (to make room for TOTAL)
    max_data_rows = 14  # Fixed to 14 rows to allow for TOTAL row
    
    logger.info(f"📏 Slide constraints: {constraints['max_table_height']:.2f}\" height available")
    logger.info(f"📏 FIXED: Max rows per slide set to {max_data_rows} data rows + 1 header")
    
    return max_data_rows

def add_slide_number(slide, slide_number, total_slides=None):
    """Add slide number to the bottom right corner of the slide"""
    try:
        # Create text box for slide number in bottom right corner
        slide_number_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(6.8),  # Bottom right position
            Inches(1.5), Inches(0.5)    # Width and height
        )
        
        slide_number_frame = slide_number_box.text_frame
        slide_number_frame.margin_left = 0
        slide_number_frame.margin_right = 0
        slide_number_frame.margin_top = 0
        slide_number_frame.margin_bottom = 0
        
        # Format slide number text - simple format like "1", "2", "3"
        slide_number_frame.text = str(slide_number)
        
        p = slide_number_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        
        logger.info(f"✅ Added slide number {slide_number} to slide")
        
    except Exception as e:
        logger.error(f"Error adding slide number: {e}")

def create_thank_you_slide(prs, logo_file=None):
    """Create a simple 'Thank You' slide as the last slide - CLEAN VERSION"""
    try:
        blank_slide_layout = prs.slide_layouts[6]
        thank_you_slide = prs.slides.add_slide(blank_slide_layout)
        
        # Main "Thank You" text - centered vertically and horizontally
        thank_you_text = thank_you_slide.shapes.add_textbox(
            Inches(0.5), Inches(3.0),  # Centered vertically
            Inches(12.33), Inches(1.5)
        )
        thank_you_frame = thank_you_text.text_frame
        thank_you_frame.text = "THANK YOU"
        p = thank_you_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        # REMOVED: Subtitle "FOR YOUR ATTENTION" - as per requirement
        
        return thank_you_slide
        
    except Exception as e:
        logger.error(f"Error creating thank you slide: {e}")
        raise

def finalize_presentation_with_numbering_and_thank_you(prs, logo_file=None):
    """Add slide numbers to all slides and create thank you slide"""
    try:
        # Create thank you slide
        create_thank_you_slide(prs, logo_file)
        
        # Get total number of slides
        total_slides = len(prs.slides)
        
        # Add slide numbers to all slides except the first (title slide)
        for i, slide in enumerate(prs.slides):
            if i == 0:  # Skip title slide
                continue
            else:  # All other slides including thank you slide - simple numbering
                add_slide_number(slide, i)  # Changed from i+1 to i for proper numbering
        
        logger.info(f"✅ Added simple slide numbers to {total_slides} slides and created thank you slide")
        
    except Exception as e:
        logger.error(f"Error finalizing presentation: {e}")
        raise

def format_column_header(header_text):
    """Format column headers with proper units and uppercase"""
    header_text = str(header_text).upper()
    
    # Enhanced label formatting for Product Growth
    if 'LY_QTY' in header_text:
        header_text = header_text.replace('LY_QTY', 'LAST-YEAR QTY/MT')
    elif 'LY_VALUE' in header_text:
        header_text = header_text.replace('LY_VALUE', 'LAST-YEAR VALUE/L')
    elif 'CY_QTY' in header_text:
        header_text = header_text.replace('CY_QTY', 'CURRENT-YEAR QTY/MT')
    elif 'CY_VALUE' in header_text:
        header_text = header_text.replace('CY_VALUE', 'CURRENT-YEAR VALUE/L')
    elif 'BUDGET_QTY' in header_text:
        header_text = header_text.replace('BUDGET_QTY', 'TARGET QTY/MT')
    elif 'BUDGET_VALUE' in header_text:
        header_text = header_text.replace('BUDGET_VALUE', 'TARGET VALUE/L')
    elif 'ACHIEVEMENT' in header_text:  # CHANGED TO GROWTH
        header_text = 'GROWTH %'
    # OD Target vs Collection specific renaming
    elif '% ACHIEVED (SELECTED MONTH)' in header_text:
        header_text = 'FOR THE MONTH % ACHIEVED'
    # General BUDGET to TARGET replacement
    elif 'BUDGET QTY' in header_text:
        header_text = header_text.replace('BUDGET QTY', 'TARGET QTY')
    elif 'BUDGET VALUE' in header_text:
        header_text = header_text.replace('BUDGET VALUE', 'TARGET VALUE')
    elif header_text == 'BUDGET':
        header_text = 'TARGET'
    # Add units to quantity and value columns
    elif any(keyword in header_text for keyword in ['QTY', 'QUANTITY']):
        if '/MT' not in header_text:
            header_text = header_text.replace('QTY', 'QTY/MT').replace('QUANTITY', 'QUANTITY/MT')
    elif any(keyword in header_text for keyword in ['VALUE', 'AMOUNT', 'COLLECTION', 'TARGET', 'BILLED']):
        if 'ACHIEVEMENT' not in header_text and '%' not in header_text and 'GROWTH' not in header_text:
            if '/L' not in header_text:
                header_text += '/L'
    
    return header_text

def format_slide_title(title):
    """Remove quantity and value labels from slide titles and format properly with enhanced Product Growth labels"""
    # FIRST: Handle Product Growth titles specifically
    title_lower = title.lower()
    
    # If it's a Product Growth title, preserve the company name and growth type
    if any(keyword in title_lower for keyword in ['quantity growth', 'value growth']):
        # Don't modify Product Growth titles - they should already be in correct format
        return title
    
    # Enhanced Product Growth title formatting (for legacy compatibility)
    title = title.replace(' LY:', ' Last-Year:').replace(' CY:', ' Current-Year:')
    title = title.replace('LY:', 'Last-Year:').replace('CY:', 'Current-Year:')
    
    # Remove common quantity and value indicators from titles (but NOT from Product Growth)
    patterns_to_remove = [
        r'\(Qty in Mt\)',
        r'\(Value in Lakhs\)',
        r'\(Qty in MT\)',
        r'\(Value in L\)',
        # REMOVED: These patterns so Product Growth titles aren't affected
        # r'Quantity Growth',
        # r'Value Growth',
        # r'- Quantity Growth',
        # r'- Value Growth'
    ]
    
    for pattern in patterns_to_remove:
        title = re.sub(pattern, '', title, flags=re.IGNORECASE)
    
    # Clean up extra spaces and dashes (but preserve Product Growth format)
    title = re.sub(r'\s+', ' ', title)
    title = re.sub(r'\s*-\s*$', '', title)
    title = title.strip()
    
    return title.upper()

def create_title_slide(prs, title, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create title slide with enhanced formatting and executive information - REMOVED 'ACCLLP' SUBTITLE"""
    try:
        blank_slide_layout = prs.slide_layouts[6]
        title_slide = prs.slides.add_slide(blank_slide_layout)
        
        # Company name
        company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        company_frame = company_name.text_frame
        company_frame.text = "ASIA CRYSTAL COMMODITY LLP"
        p = company_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Calibri"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        
        # Logo (if provided)
        if logo_file is not None:
            try:
                if isinstance(logo_file, str):
                    # Base64 encoded logo
                    logo_data = base64.b64decode(logo_file)
                    logo_buffer = BytesIO(logo_data)
                    title_slide.shapes.add_picture(
                        logo_buffer, Inches(5.665), Inches(1.5), 
                        width=Inches(2), height=Inches(2)
                    )
                else:
                    # File object
                    logo_buffer = BytesIO(logo_file.read())
                    title_slide.shapes.add_picture(
                        logo_buffer, Inches(5.665), Inches(1.5), 
                        width=Inches(2), height=Inches(2)
                    )
                    logo_file.seek(0)  # Reset for reuse
            except Exception as e:
                logger.error(f"Error adding logo to slide: {str(e)}")
        
        # Title
        formatted_title = format_slide_title(title)
        title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = formatted_title
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Calibri"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        # REMOVED: Subtitle with "ACCLLP" - as per requirement
        
        # Executive Information Section - moved up since subtitle is removed
        info_y_position = 5.0
        
        # Executive Name
        if executive_name:
            name_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(info_y_position), Inches(12.33), Inches(0.5))
            name_frame = name_box.text_frame
            name_frame.text = f"NAME: {executive_name.upper()}"
            p = name_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Calibri"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(128, 0, 128)
            info_y_position += 0.4
        
        # Date
        if date_str:
            date_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(info_y_position), Inches(12.33), Inches(0.5))
            date_frame = date_box.text_frame
            date_frame.text = f"DATE: {date_str}"
            p = date_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Calibri"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(128, 0, 128)
            info_y_position += 0.4
        
        # Branch Name
        if branch_name:
            branch_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(info_y_position), Inches(12.33), Inches(0.5))
            branch_frame = branch_box.text_frame
            branch_frame.text = f"BRANCH: {branch_name.upper()}"
            p = branch_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Calibri"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(128, 0, 128)
        
        return title_slide
        
    except Exception as e:
        logger.error(f"Error creating title slide: {e}")
        raise

def is_total_row(row_value):
    if pd.isna(row_value):
        return False
    
    first_col_value = str(row_value).strip().upper()
    
    total_patterns = [
        'TOTAL', 'GRAND TOTAL', 'OVERALL TOTAL', 'OVERALL SALES',
        'OVERALL SALES TOTAL', 'SALES TOTAL', 'SUPER', 'SUPER TOTAL',
        'SUBTOTAL', 'SUB TOTAL', 'SUMMARY', 'OD COLLECTION',
        'OD COLLECTION TOTAL', 'OD TARGET', 'OD TARGET TOTAL'
    ]
    
    # Exact match
    if first_col_value in total_patterns:
        return True
    
    # Partial match
    for pattern in total_patterns:
        if pattern in first_col_value:
            return True
    
    # Special cases
    if first_col_value.endswith('TOTAL') or first_col_value.startswith('SUPER'):
        return True
    
    return False


def format_slide_title_custom(title: str) -> str:
    """
    Format slide title:
    1️⃣ Convert to uppercase
    2️⃣ Replace 'FY YYYY-YYYY' with just first year
    """
    import re
    # Replace 'FY 2025-2026' or similar with '2025'
    title = re.sub(r'FY (\d{4})-\d{4}', r'\1', title, flags=re.IGNORECASE)
    # Convert to uppercase
    title = title.upper()
    return title

def format_customer_title_with_month(title, sorted_months=None, df=None):
    """
    Format customer report title to include month in the format:
    'NUMBER OF BILLED CUSTOMERS - JUL 2025'
    Extract month from actual data columns if df is provided
    """
    # Convert to uppercase
    title = title.upper()
    
    # If it's a customer report, ensure proper month formatting
    if "NUMBER OF BILLED CUSTOMERS" in title:
        
        # PRIORITY 1: Extract month from DataFrame columns (most accurate)
        if df is not None and not df.empty:
            # Look for month columns in the DataFrame
            month_columns = []
            for col in df.columns:
                col_str = str(col).upper()
                # Check if column contains month and year pattern like "JUL 2025"
                import re
                month_year_pattern = r'([A-Z]{3})\s+(\d{4})'
                if re.search(month_year_pattern, col_str):
                    month_columns.append(col_str)
            
            # Use the first month column found
            if month_columns:
                month_col = month_columns[0]
                match = re.search(month_year_pattern, month_col)
                if match:
                    month_part = match.group(1)  # e.g., "JUL"
                    year_part = match.group(2)   # e.g., "2025"
                    return f"NUMBER OF BILLED CUSTOMERS - {month_part} {year_part}"
        
        # PRIORITY 2: Use sorted_months if provided
        if sorted_months and len(sorted_months) > 0:
            # Get the most recent month
            recent_month = sorted_months[-1] if isinstance(sorted_months, list) else sorted_months
            
            # Extract month and year from the recent_month string
            if isinstance(recent_month, str):
                recent_month_upper = recent_month.upper()
                import re
                # Look for month-year pattern in sorted_months
                month_year_match = re.search(r'([A-Z]{3})\s+(\d{4})', recent_month_upper)
                if month_year_match:
                    month_part = month_year_match.group(1)
                    year_part = month_year_match.group(2)
                    return f"NUMBER OF BILLED CUSTOMERS - {month_part} {year_part}"
                elif " " in recent_month:
                    # Format like "JUL 2025" 
                    parts = recent_month_upper.split()
                    if len(parts) >= 2:
                        month_part = parts[0][:3]  # First 3 chars of first part
                        year_part = parts[-1]      # Last part as year
                        return f"NUMBER OF BILLED CUSTOMERS - {month_part} {year_part}"
        
        # PRIORITY 3: Fallback - extract year from title and use current month
        import re
        year_match = re.search(r'\b(\d{4})\b', title)
        if year_match:
            year = year_match.group(1)
            import datetime
            current_month = datetime.datetime.now().strftime("%b").upper()
            return f"NUMBER OF BILLED CUSTOMERS - {current_month} {year}"
    
    return title

def fix_budget_vs_billed_column_order(df, report_type="auto"):
    """
    ENHANCED: Pre-process DataFrame to ensure correct column ordering for ALL Target vs Billed reports
    Fixed to handle both "Target Against Billed" and "Overall Sales" reports properly
    
    Args:
        df: DataFrame with budget vs billed data
        report_type: "qty", "value", or "auto" to detect automatically
    
    Returns:
        DataFrame with correctly ordered columns: Executive → Target → Billed [→ %]
    """
    try:
        if df.empty:
            return df
            
        available_columns = list(df.columns)
        logger.info(f"🔍 Original columns: {available_columns}")
        
        # Auto-detect report type if not specified
        if report_type == "auto":
            # More comprehensive detection
            has_qty = any(keyword in col.lower() for col in available_columns 
                         for keyword in ['qty', 'quantity', 'mt'])
            has_value = any(keyword in col.lower() for col in available_columns 
                           for keyword in ['value', 'lakh', 'amount'])
            
            if has_qty and not has_value:
                report_type = "qty"
            elif has_value and not has_qty:
                report_type = "value"
            elif has_qty and has_value:
                # Both exist, determine primary type by counting occurrences
                qty_count = sum(1 for col in available_columns 
                               if any(kw in col.lower() for kw in ['qty', 'quantity', 'mt']))
                value_count = sum(1 for col in available_columns 
                                 if any(kw in col.lower() for kw in ['value', 'lakh', 'amount']))
                
                report_type = "qty" if qty_count >= value_count else "value"
            else:
                logger.warning("Cannot determine report type, returning original order")
                return df
        
        # FIXED: More flexible and accurate column matching
        def find_column_by_patterns(patterns, available_cols):
            """Find column that matches any of the patterns (case insensitive)"""
            for col in available_cols:
                col_lower = col.lower()
                for pattern in patterns:
                    if pattern.lower() in col_lower:
                        return col
            return None
        
        # Find the actual column names in the DataFrame
        executive_col = find_column_by_patterns(['executive'], available_columns)
        
        if report_type == "qty":
            # FIXED: More comprehensive patterns for quantity columns
            target_patterns = [
                'target qty', 'budget qty', 'target quantity', 'budget quantity',
                'target', 'budget'  # Fallback patterns
            ]
            billed_patterns = [
                'billed qty', 'billed quantity', 'actual qty', 'actual quantity',
                'billed'  # Fallback pattern
            ]
            
            # Find target and billed columns with priority
            target_col = None
            for pattern in target_patterns:
                target_col = find_column_by_patterns([pattern], available_columns)
                if target_col:
                    break
            
            billed_col = None
            for pattern in billed_patterns:
                billed_col = find_column_by_patterns([pattern], available_columns)
                if billed_col and billed_col != target_col:  # Ensure it's not the same column
                    break
            
        elif report_type == "value":
            # FIXED: More comprehensive patterns for value columns  
            target_patterns = [
                'target value', 'budget value', 'target amount', 'budget amount',
                'target', 'budget'  # Fallback patterns
            ]
            billed_patterns = [
                'billed value', 'billed amount', 'actual value', 'actual amount',
                'billed'  # Fallback pattern
            ]
            
            # Find target and billed columns with priority
            target_col = None
            for pattern in target_patterns:
                target_col = find_column_by_patterns([pattern], available_columns)
                if target_col:
                    break
            
            billed_col = None
            for pattern in billed_patterns:
                billed_col = find_column_by_patterns([pattern], available_columns)
                if billed_col and billed_col != target_col:  # Ensure it's not the same column
                    break
        else:
            return df
        
        # Find percentage column
        percentage_col = find_column_by_patterns(['%', 'percent', 'achievement'], available_columns)
        
        # Build the CORRECT order: Executive → Target → Billed [→ %]
        correct_order = []
        if executive_col:
            correct_order.append(executive_col)
        if target_col:
            correct_order.append(target_col)
        if billed_col:
            correct_order.append(billed_col)
        if percentage_col:
            correct_order.append(percentage_col)
        
        # Add any remaining columns that weren't matched
        for col in available_columns:
            if col not in correct_order:
                correct_order.append(col)
        
        # Only reorder if we found the expected columns
        if len(correct_order) >= 3 and executive_col and target_col and billed_col:
            # Return reordered DataFrame
            reordered_df = df[correct_order].copy()
            logger.info(f"✅ Fixed column order: {available_columns} → {correct_order}")
            return reordered_df
        else:
            logger.warning(f"Could not identify all required columns (Executive: {executive_col}, Target: {target_col}, Billed: {billed_col}), keeping original order")
            return df
        
    except Exception as e:
        logger.error(f"Error fixing column order: {e}")
        return df

def add_table_slide(prs, df, title, percent_cols=None, is_consolidated=False, max_rows=15, sorted_months=None):
    """Add table slide with enforced TOTAL row styling, guaranteed gray background for TOTAL rows, and FIXED column ordering."""
    try:
        if percent_cols is None:
            percent_cols = []

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Detect if it's a customer report to remove S.No column
        title_lower = title.lower()
        is_customer_report = any(phrase in title_lower for phrase in [
            "number of billed customers", 
            "billed customers", 
            "customer analysis",
            "nbc"
        ])

        # Remove variations of S.No column for customer reports
        if is_customer_report:
            sno_variations = ['S.No', 'SNo', 'S No', 'sno', 's.no', 'Serial No', 'Sr. No']
            for sno_col in sno_variations:
                if sno_col in df.columns:
                    df = df.drop(sno_col, axis=1)
                    logger.info(f"✅ Removed {sno_col} column from customer report: {title}")

        # ENHANCED: Apply column ordering fix for target vs billed and overall sales reports
        if any(phrase in title_lower for phrase in [
            "target against billed", "target vs billed", "target v/s billed", 
            "overall sales", "overall sale", "total sales"
        ]):
            # Determine report type from title
            if any(word in title_lower for word in ["quantity", "qty", "mt"]):
                df = fix_budget_vs_billed_column_order(df, "qty")
            elif any(word in title_lower for word in ["value", "lakhs", "amount"]):
                df = fix_budget_vs_billed_column_order(df, "value")
            else:
                # Auto-detect from columns
                df = fix_budget_vs_billed_column_order(df, "auto")

        # Format title
        if is_customer_report:
            formatted_title = format_customer_title_with_month(title, sorted_months, df)
        else:
            formatted_title = format_slide_title_custom(title)
            
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = formatted_title
        p = title_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.CENTER

        logger.info(f"🔍 PROCESSING SLIDE: '{formatted_title}' with max {max_rows} rows")
        logger.info(f"🔍 DataFrame columns: {list(df.columns)}")
        logger.info(f"🔍 DataFrame shape: {df.shape}")

        actual_columns = list(df.columns)
        title_lower = title.lower()

        # Enhanced column ordering logic with FIXED ordering for target vs billed reports
        ordered_columns = []
        strict_order = df.attrs.get("columns", [])
        
        if not strict_order:
            if any(phrase in title_lower for phrase in ["target against billed", "target vs billed", "target v/s billed"]):
                # For Target Against Billed reports, use the fixed order from DataFrame (after column fixing)
                ordered_columns = actual_columns.copy()
                
            elif any(phrase in title_lower for phrase in ["overall sales", "overall sale", "total sales"]):
                # For Overall Sales reports, use the fixed order from DataFrame (after column fixing)
                ordered_columns = actual_columns.copy()
                
            elif any(phrase in title_lower for phrase in ["od target vs collection", "od collection"]):
                strict_order = [
                    'Executive', 'Due Target', 'Collection Achieved',
                    'Overall % Achieved', 'For the month Overdue',
                    'For the month Collection', '% Achieved (Selected Month)'
                ]
                # Apply strict order
                for col in strict_order:
                    if col in actual_columns:
                        ordered_columns.append(col)
                for col in actual_columns:
                    if col not in ordered_columns:
                        ordered_columns.append(col)
            else:
                ordered_columns = actual_columns.copy()
        else:
            # Use frontend-defined order
            for col in strict_order:
                if col in actual_columns:
                    ordered_columns.append(col)
            for col in actual_columns:
                if col not in ordered_columns:
                    ordered_columns.append(col)

        # If ordered_columns is still empty, use actual_columns
        if not ordered_columns:
            ordered_columns = actual_columns.copy()

        # Table creation
        num_rows = len(df) + 1
        num_cols = len(ordered_columns)
        table_width = Inches(12.5)
        table_height = Inches(len(df) * 0.35 + 0.5)
        table_left = Inches(0.4)
        table_top = Inches(1.3)

        table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height).table

        # Enhanced column width adjustments
        if num_cols == 7 and any(phrase in title_lower for phrase in ["od target vs collection", "od collection"]):
            table.columns[0].width = Inches(2.8)
            table.columns[1].width = Inches(1.6)
            table.columns[2].width = Inches(1.6)
            table.columns[3].width = Inches(1.3)
            table.columns[4].width = Inches(1.6)
            table.columns[5].width = Inches(1.9)
            table.columns[6].width = Inches(1.7)
        elif num_cols > 0:
            if num_cols == 2:
                table.columns[0].width = Inches(7.0)
                table.columns[1].width = Inches(5.5)
            elif is_customer_report and ('Executive' in ordered_columns[0] or 'Executive Name' in ordered_columns[0]):
                table.columns[0].width = Inches(4.5)
                remaining_width = 8.0
                if num_cols > 1:
                    col_width = remaining_width / (num_cols - 1)
                    for i in range(1, num_cols):
                        table.columns[i].width = Inches(col_width)
            elif 'Executive' in ordered_columns[0]:
                table.columns[0].width = Inches(4.0)
                remaining_width = 8.5
                if num_cols > 1:
                    col_width = remaining_width / (num_cols - 1)
                    for i in range(1, num_cols):
                        table.columns[i].width = Inches(col_width)
            else:
                col_width = 12.5 / num_cols
                for i in range(num_cols):
                    table.columns[i].width = Inches(col_width)

        # Header formatting
        for i, col_name in enumerate(ordered_columns):
            cell = table.cell(0, i)
            
            # Enhanced header formatting for target vs billed reports
            display_header = col_name
            if any(phrase in title_lower for phrase in ["target against billed", "overall sales"]):
                if 'Budget Value' in col_name:
                    display_header = 'TARGET VALUE/L'
                elif 'Budget Qty' in col_name:
                    display_header = 'TARGET QTY/MT'
                elif 'Billed Value' in col_name:
                    display_header = 'BILLED VALUE/L'
                elif 'Billed Qty' in col_name:
                    display_header = 'BILLED QTY/MT'
                else:
                    display_header = format_column_header(col_name)
            else:
                display_header = format_column_header(col_name)
            
            cell.text = display_header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Calibri"
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)

        # Data rows formatting with guaranteed TOTAL detection and styling
        for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
            # Enhanced TOTAL row detection - check all columns for TOTAL indicators
            is_total = any(
                is_total_row(str(val)) 
                for val in row.values if pd.notna(val)
            )

            for col_idx, col_name in enumerate(ordered_columns):
                cell = table.cell(row_idx, col_idx)
                value = row.get(col_name, "")

                # Handle percentage columns
                if col_name in ['ACHIEVEMENT %', 'GROWTH %', '%', 'Overall % Achieved', '% Achieved (Selected Month)']:
                    if isinstance(value, (int, float)) and not pd.isna(value):
                        cell.text = f"{value:.2f}%"
                    else:
                        cell.text = str(value) if value is not None else ""
                elif col_idx in percent_cols and isinstance(value, (int, float)) and not pd.isna(value):
                    cell.text = f"{value:.2f}%"
                elif isinstance(value, (int, float)) and not pd.isna(value):
                    cell.text = str(int(round(value))) if abs(value - round(value)) < 0.001 else f"{value:.2f}"
                else:
                    cell.text = str(value).upper() if value is not None and not pd.isna(value) else ""

                para = cell.text_frame.paragraphs[0]
                para.font.name = "Calibri"
                para.alignment = PP_ALIGN.CENTER
                para.font.size = Pt(14)
                cell.margin_left = Inches(0.03)
                cell.margin_right = Inches(0.03)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                cell.fill.solid()

                # GUARANTEED TOTAL ROW STYLING with gray background
                if is_total:
                    # Apply gray background and bold formatting for TOTAL rows
                    cell.fill.fore_color.rgb = RGBColor(169, 169, 169)  # Gray background
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(0, 0, 0)  # Black text for readability
                    logger.info(f"✅ Applied TOTAL row styling to row {row_idx}, col {col_idx}: '{cell.text}'")
                else:
                    # Regular row styling with alternating colors
                    para.font.bold = False
                    para.font.color.rgb = RGBColor(0, 0, 0)
                    # Apply alternating background colors
                    if row_idx % 2 == 0:
                        cell.fill.fore_color.rgb = RGBColor(221, 235, 247)  # Light blue
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

        logger.info(f"✅ Created table slide: '{formatted_title}' with {len(df)} rows and {num_cols} columns")
        return slide

    except Exception as e:
        logger.error(f"Error adding table slide: {e}")
        raise

def process_df_for_slides(prs, df, title_base, percent_cols=None, is_consolidated=False):
    """Process DataFrame for slides with FIXED 14 rows per slide limit - ENSURE ALL DATA IS INCLUDED"""
    try:
        if percent_cols is None:
            percent_cols = []
        
        # Conditionally remove non-TOTAL "ACCLP" rows to preserve potential TOTAL rows
        if not df.empty:
            first_col = df.columns[0]
            df = df[~((df[first_col] == "ACCLP") & (~df[first_col].apply(is_total_row)))].copy()
        
        # Separate data rows and total row - ENHANCED DETECTION
        first_col = df.columns[0]
        total_mask = df[first_col].apply(is_total_row)
        data_rows = df[~total_mask].copy()
        total_row = df[total_mask].copy()
        
        num_data_rows = len(data_rows)
        
        # FIXED: Use 14 data rows per slide for ALL reports to make room for TOTAL
        split_threshold = 14
        
        logger.info(f"🔄 Processing {num_data_rows} data rows + {len(total_row)} total rows with threshold {split_threshold}")
        print(f"🔄 Processing {num_data_rows} data rows + {len(total_row)} total rows")
        
        if num_data_rows == 0:
            # Only TOTAL row(s)
            if not total_row.empty:
                add_table_slide(prs, total_row, title_base, percent_cols=percent_cols, is_consolidated=is_consolidated, max_rows=15)
            return
        
        if num_data_rows <= split_threshold:
            # Single slide - add total row at the end if it exists
            if not total_row.empty:
                combined_df = pd.concat([data_rows, total_row], ignore_index=True)
            else:
                combined_df = data_rows.copy()
            add_table_slide(prs, combined_df, title_base, percent_cols=percent_cols, is_consolidated=is_consolidated, max_rows=15)
            return
        
        # Split logic for more data rows - ENSURE ALL DATA IS INCLUDED
        num_parts = math.ceil(num_data_rows / split_threshold)
        logger.info(f"📊 Splitting {num_data_rows} data rows into {num_parts} parts")
        print(f"📊 Splitting {num_data_rows} data rows into {num_parts} parts")
        
        # Create slides for each part - ENSURE ALL DATA IS INCLUDED
        for i in range(num_parts):
            start_idx = i * split_threshold
            end_idx = min((i + 1) * split_threshold, num_data_rows)
            part = data_rows.iloc[start_idx:end_idx].copy()
            
            # Check if this is the LAST part
            is_last_part = (i == num_parts - 1)
            
            # If this is the last part AND we have a total row, add it
            if is_last_part and not total_row.empty:
                combined_part = pd.concat([part, total_row], ignore_index=True)
            else:
                combined_part = part.copy()
            
            # Ensure all numeric columns are properly rounded
            numeric_cols = combined_part.select_dtypes(include=[np.number]).columns
            for col in numeric_cols:
                combined_part[col] = combined_part[col].apply(lambda x: round(float(x), 2) if pd.notna(x) else x)
            
            # Use the original title_base without appending part numbers
            part_title = title_base  # Changed from f"{title_base} - Part {i+1}"
            add_table_slide(prs, combined_part, part_title, percent_cols=percent_cols, is_consolidated=is_consolidated, max_rows=15)
            
            if is_last_part and not total_row.empty:
                logger.info(f"✅ Created slide with {len(part)} data rows + TOTAL ROW")
            else:
                logger.info(f"✅ Created slide with {len(part)} data rows")
        
    except Exception as e:
        logger.error(f"Error processing DataFrame for slides: {e}")
        raise


def create_consolidated_ppt(dfs_info, logo_file=None, title="Consolidated Report", executive_name=None, date_str=None, branch_name=None):
    """Create consolidated PowerPoint with enhanced formatting, slide numbers and thank you slide - ENSURE ALL DATA INCLUDED"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide with executive information
        create_title_slide(prs, title, logo_file, executive_name, date_str, branch_name)
        
        logger.info(f"📊 Creating consolidated PPT with {len(dfs_info)} reports (ALL DATA INCLUDED)")
        print(f"📊 Creating consolidated PPT with {len(dfs_info)} reports (ALL DATA INCLUDED)")
        
        # Process each report with enhanced formatting - ENSURE ALL DATA IS INCLUDED
        for df_info in dfs_info:
            df_data = df_info.get('df', [])
            slide_title = df_info['title']
            percent_cols = df_info.get('percent_cols', [])
            
            # ENHANCED: Fix Product Growth titles for consolidated PPT
            original_lower = slide_title.lower()
            
            # Pattern 1: "CompanyName - Quantity Growth - Extra Text" - Keep as Quantity Growth
            if 'quantity growth' in original_lower:
                # Extract company name (first part before any " - ")
                company_name = slide_title.split(' - ')[0].strip()
                slide_title = f"{company_name} - Quantity Growth"
                logger.info(f"✅ FIXED Quantity Growth title: '{slide_title}'")
            
            # Pattern 2: "CompanyName - Value Growth - Extra Text" 
            elif 'value growth' in original_lower:
                company_name = slide_title.split(' - ')[0].strip()
                slide_title = f"{company_name} - Value Growth"
                logger.info(f"✅ FIXED Value Growth title: '{slide_title}'")
            
            # Pattern 3: Handle malformed titles like "Cp - - Last Year..."
            elif ' - - ' in slide_title and ('last year' in original_lower or 'current year' in original_lower):
                company_name = slide_title.split(' - ')[0].strip()
                # Check if this might be quantity or value based on column headers
                columns_str = str(df_info.get('columns', [])).lower()
                if 'qty' in columns_str or 'quantity' in columns_str:
                    slide_title = f"{company_name} - Quantity Growth"
                else:
                    slide_title = f"{company_name} - Value Growth"
                logger.info(f"✅ FIXED malformed title: '{slide_title}'")
            
            # Convert to DataFrame if it's a list of dicts
            if isinstance(df_data, list):
                df = pd.DataFrame(df_data)
            else:
                df = df_data.copy()

            # Apply frontend-defined column order (if present)
            frontend_order = df_info.get("columns")
            if frontend_order:
                df = df[[col for col in frontend_order if col in df.columns]]
                logger.info(f"📐 Applied frontend column order: {frontend_order}")
            
            if df.empty:
                logger.warning(f"Skipping empty report: {slide_title}")
                continue
            
            logger.info(f"🔄 Processing consolidated report: {slide_title}")
            print(f"🔄 Processing consolidated report: {slide_title}")
            print(f"🔍 Original DataFrame columns: {list(df.columns)}")
            
            # Use process_df_for_slides for proper splitting
            process_df_for_slides(prs, df, slide_title, percent_cols=percent_cols, is_consolidated=True)
        
        # Finalize with slide numbers and thank you slide
        finalize_presentation_with_numbering_and_thank_you(prs, logo_file)
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        # Log final slide count
        actual_slide_count = len(prs.slides)
        logger.info(f"📊 Consolidated PPT created with {actual_slide_count} total slides (ALL DATA INCLUDED)")
        print(f"📊 Consolidated PPT created with {actual_slide_count} total slides (ALL DATA INCLUDED)")
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating consolidated PPT: {str(e)}")
        raise Exception(f"Error creating consolidated PPT: {str(e)}")

def create_product_growth_ppt(group_results, month_title, logo_file=None, ly_month=None, cy_month=None, executive_name=None, date_str=None, branch_name=None):
    """
    Create Product Growth PPT with FIXED column ordering and enhanced formatting - ENSURE TOTAL ROW IS INCLUDED
    Updated to use "GROWTH" instead of "ACHIEVEMENT"
    """
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Enhanced title slide text generation - ENSURE UPPERCASE
        if ly_month and cy_month:
            title_slide_text = f"PRODUCT GROWTH – LAST-YEAR: {ly_month.upper()} VS CURRENT-YEAR: {cy_month.upper()}"
            logger.info(f"✅ Title with months: {title_slide_text}")
        elif month_title and month_title != "Product Growth Analysis":
            title_slide_text = f"PRODUCT GROWTH – {month_title.upper()}"
            logger.info(f"✅ Title with custom title: {title_slide_text}")
        else:
            title_slide_text = "PRODUCT GROWTH ANALYSIS"
            logger.info(f"✅ Default title: {title_slide_text}")
        
        print(f"🔍 Creating Product Growth PPT: {title_slide_text}")
        
        # Create title slide
        create_title_slide(prs, title_slide_text, logo_file, executive_name, date_str, branch_name)
        
        def fix_product_growth_columns(df, report_type="auto"):
            """Fix column ordering for Product Growth DataFrames and ensure TOTAL row is present"""
            try:
                if df.empty:
                    return df
                    
                available_columns = list(df.columns)
                logger.info(f"🔧 Fixing columns for {report_type}: {available_columns}")
                
                # Check if TOTAL row exists
                has_total = False
                if 'PRODUCT GROUP' in available_columns:
                    has_total = df['PRODUCT GROUP'].apply(is_total_row).any()
                    logger.info(f"🔍 TOTAL row found: {has_total}")
                
                # Auto-detect report type if not specified
                if report_type == "auto":
                    has_qty = any('QTY' in col for col in available_columns)
                    has_value = any('VALUE' in col for col in available_columns)
                    
                    if has_qty and not has_value:
                        report_type = "qty"
                    elif has_value and not has_qty:
                        report_type = "value"
                    else:
                        # Mixed or unclear, count occurrences
                        qty_count = sum(1 for col in available_columns if 'QTY' in col)
                        value_count = sum(1 for col in available_columns if 'VALUE' in col)
                        
                        if qty_count > value_count:
                            report_type = "qty"
                        elif value_count > qty_count:
                            report_type = "value"
                        else:
                            # Can't determine, return as-is
                            logger.warning(f"⚠️ Cannot determine report type, using original order")
                            return df
                
                # Define correct column order
                if report_type == "qty":
                    correct_order = ['PRODUCT GROUP', 'LY_QTY', 'BUDGET_QTY', 'CY_QTY', 'ACHIEVEMENT %']
                elif report_type == "value":
                    correct_order = ['PRODUCT GROUP', 'LY_VALUE', 'BUDGET_VALUE', 'CY_VALUE', 'ACHIEVEMENT %']
                else:
                    return df
                
                # Reorder columns
                ordered_columns = []
                for col in correct_order:
                    if col in available_columns:
                        ordered_columns.append(col)
                
                # Add any remaining columns
                for col in available_columns:
                    if col not in ordered_columns:
                        ordered_columns.append(col)
                
                # Return reordered DataFrame
                reordered_df = df[ordered_columns].copy()
                logger.info(f"✅ Fixed {report_type} columns: {available_columns} → {ordered_columns}")
                return reordered_df
                
            except Exception as e:
                logger.error(f"Error fixing columns: {e}")
                return df
        
        # Process each company
        for company, data in group_results.items():
            logger.info(f"🔄 Processing company: {company}")
            
            # Handle both DataFrame and dict formats
            if isinstance(data['qty_df'], list):
                qty_df = pd.DataFrame(data['qty_df'])
            else:
                qty_df = data['qty_df'].copy()
                
            if isinstance(data['value_df'], list):
                value_df = pd.DataFrame(data['value_df'])
            else:
                value_df = data['value_df'].copy()
            
            # FIX COLUMN ORDERING
            qty_df = fix_product_growth_columns(qty_df, "qty")
            value_df = fix_product_growth_columns(value_df, "value")
            
            # Round numeric columns
            numeric_cols_qty = qty_df.select_dtypes(include=[np.number]).columns
            for col in numeric_cols_qty:
                qty_df[col] = qty_df[col].round(2)
            
            numeric_cols_value = value_df.select_dtypes(include=[np.number]).columns
            for col in numeric_cols_value:
                value_df[col] = value_df[col].round(2)
            
            # Create slide titles - ENSURE UPPERCASE
            if ly_month and cy_month:
                qty_title = f"{company.upper()} - QUANTITY GROWTH"
                value_title = f"{company.upper()} - VALUE GROWTH"
            else:
                qty_title = f"{company.upper()} - QUANTITY GROWTH"
                value_title = f"{company.upper()} - VALUE GROWTH"
            
            # FIXED: Use process_df_for_slides to ensure TOTAL row appears properly
            process_df_for_slides(prs, qty_df, qty_title, percent_cols=[4], is_consolidated=False)  # ACHIEVEMENT % column
            process_df_for_slides(prs, value_df, value_title, percent_cols=[4], is_consolidated=False)  # ACHIEVEMENT % column
        
        # Add slide numbers and thank you slide
        finalize_presentation_with_numbering_and_thank_you(prs, logo_file)
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        # Log final summary
        actual_slides = len(prs.slides)
        company_count = len(group_results)
        logger.info(f"📊 Product Growth PPT created: {company_count} companies → {actual_slides} slides (TOTAL ROWS INCLUDED)")
        print(f"✅ Product Growth PPT created: {company_count} companies → {actual_slides} slides (TOTAL ROWS INCLUDED)")
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating Product Growth PPT: {e}")
        print(f"❌ Error creating Product Growth PPT: {e}")
        return None

def create_executive_budget_ppt(results_data, month_title=None, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create PowerPoint presentation for executive target vs billed with FIXED column ordering for ALL slides"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide with executive information
        title = f"Monthly Review Meeting – {month_title}" if month_title else "Executive Target vs Billed Analysis"
        create_title_slide(prs, title, logo_file, executive_name, date_str, branch_name)
        
        # Convert results data to DataFrames and ENSURE ALL DATA IS PROCESSED
        budget_vs_billed_qty_df = pd.DataFrame(results_data.get('budget_vs_billed_qty', []))
        budget_vs_billed_value_df = pd.DataFrame(results_data.get('budget_vs_billed_value', []))
        overall_sales_qty_df = pd.DataFrame(results_data.get('overall_sales_qty', []))
        overall_sales_value_df = pd.DataFrame(results_data.get('overall_sales_value', []))
        
        # FIXED: Apply column ordering to ALL DataFrames with correct detection
        if not budget_vs_billed_qty_df.empty:
            budget_vs_billed_qty_df = fix_budget_vs_billed_column_order(budget_vs_billed_qty_df, "qty")
            qty_title = f"TARGET AGAINST BILLED - QUANTITY - {month_title}" if month_title else "TARGET AGAINST BILLED - QUANTITY"
            process_df_for_slides(prs, budget_vs_billed_qty_df, qty_title, percent_cols=[3], is_consolidated=False)
        
        if not budget_vs_billed_value_df.empty:
            budget_vs_billed_value_df = fix_budget_vs_billed_column_order(budget_vs_billed_value_df, "value")
            value_title = f"TARGET AGAINST BILLED - VALUE - {month_title}" if month_title else "TARGET AGAINST BILLED - VALUE"
            process_df_for_slides(prs, budget_vs_billed_value_df, value_title, percent_cols=[3], is_consolidated=False)
        
        # FIXED: Apply column ordering to Overall Sales DataFrames with proper type detection
        if not overall_sales_qty_df.empty:
            overall_sales_qty_df = fix_budget_vs_billed_column_order(overall_sales_qty_df, "qty")
            sales_qty_title = f"OVERALL SALES - QUANTITY - {month_title}" if month_title else "OVERALL SALES - QUANTITY"
            process_df_for_slides(prs, overall_sales_qty_df, sales_qty_title, percent_cols=[], is_consolidated=False)
        
        if not overall_sales_value_df.empty:
            overall_sales_value_df = fix_budget_vs_billed_column_order(overall_sales_value_df, "value")
            sales_value_title = f"OVERALL SALES - VALUE - {month_title}" if month_title else "OVERALL SALES - VALUE"
            process_df_for_slides(prs, overall_sales_value_df, sales_value_title, percent_cols=[], is_consolidated=False)
        
        # Finalize with slide numbers and thank you slide
        finalize_presentation_with_numbering_and_thank_you(prs, logo_file)
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        logger.info(f"📊 Executive Target PPT created with FIXED column ordering")
        print(f"📊 Executive Target PPT created with FIXED column ordering")
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating Executive Target PPT: {str(e)}")
        raise Exception(f"Error creating Executive Target PPT: {str(e)}")



def create_executive_od_ppt(results_data, month_title=None, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create PowerPoint presentation for executive OD vs Collection with slide numbers and thank you slide"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide with executive information
        title = f"OD Target vs Collection – {month_title}" if month_title else "OD Target vs Collection Analysis"
        create_title_slide(prs, title, logo_file, executive_name, date_str, branch_name)
        
        # Process OD results
        if 'od_results' in results_data:
            od_data = results_data['od_results']
        elif isinstance(results_data, list):
            od_data = results_data
        else:
            od_data = results_data
            
        if od_data:
            # Convert to DataFrame
            if isinstance(od_data, list):
                od_df = pd.DataFrame(od_data)
            else:
                od_df = pd.DataFrame(od_data) if not isinstance(od_data, pd.DataFrame) else od_data.copy()
            
            logger.info(f"OD DataFrame columns: {list(od_df.columns)}")
            logger.info(f"OD DataFrame shape: {od_df.shape}")
            print(f"🔍 OD DataFrame columns: {list(od_df.columns)}")
            
            # Identify percentage columns
            percentage_column_names = ['Overall % Achieved', '% Achieved (Selected Month)']
            percent_cols = []
            for i, col in enumerate(od_df.columns):
                if col in percentage_column_names:
                    percent_cols.append(i)
            
            logger.info(f"OD percentage columns identified: {percentage_column_names} at indices {percent_cols}")
            
            # Use process_df_for_slides for proper splitting with updated total row styling
            process_df_for_slides(prs, od_df, title, percent_cols=percent_cols, is_consolidated=False)
        
        # Finalize with slide numbers and thank you slide
        finalize_presentation_with_numbering_and_thank_you(prs, logo_file)
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"Error creating Executive OD PPT: {str(e)}")
        raise Exception(f"Error creating Executive OD PPT: {str(e)}")
    
def create_nbc_individual_ppt(customer_df, customer_title, sorted_months, financial_year, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create individual PPT for NBC report with slide numbers and thank you slide - ENSURE ALL DATA INCLUDED - REMOVE S.No COLUMN"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Remove S.No column if it exists
        if 'S.No' in customer_df.columns:
            customer_df = customer_df.drop('S.No', axis=1)
            logger.info("✅ Removed S.No column from customer DataFrame")
        
        # Extract month from sorted_months for title formatting
        month_only = ""
        if sorted_months and len(sorted_months) > 0:
            # Get the most recent month (assuming sorted_months is chronological)
            recent_month = sorted_months[-1] if isinstance(sorted_months, list) else sorted_months
            # Format as "NUMBER OF BILLED CUSTOMERS - JUL 2025" (just month, no extra year)
            month_only = f" - {recent_month.upper()}"
        
        # Create title slide with executive information
        title = f"NUMBER OF BILLED CUSTOMERS{month_only}"
        create_title_slide(prs, title, logo_file, executive_name, date_str, branch_name)
        
        # FIXED: Use process_df_for_slides for proper splitting of customer data
        if not customer_df.empty:
            # Format the slide title with month only
            slide_title = f"NUMBER OF BILLED CUSTOMERS{month_only}"
            process_df_for_slides(prs, customer_df, slide_title, percent_cols=[], is_consolidated=False)
        
        # Finalize with slide numbers and thank you slide
        finalize_presentation_with_numbering_and_thank_you(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        logger.info(f"✅ NBC PPT created with ALL DATA INCLUDED - Title: {title} - S.No column removed")
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating NBC PPT: {e}")
        return None

def create_customer_ppt_slide(slide, df, title, sorted_months, is_last_slide=False):
    """Create customer PPT slide with enhanced formatting and FIXED TOTAL ROW STYLING - REMOVE S.No COLUMN - INCREASED WIDTHS"""
    try:
        # Remove S.No column if it exists
        if 'S.No' in df.columns:
            df = df.drop('S.No', axis=1)
            logger.info("✅ Removed S.No column from slide DataFrame")
        
        if df.empty or len(df.columns) < 1:  # Changed from < 2 to < 1 since we removed S.No
            logger.warning(f"Skipping customer slide: DataFrame is empty or has insufficient columns {df.columns.tolist()}")
            return
        
        # Extract month from sorted_months for title formatting
        month_only = ""
        if sorted_months and len(sorted_months) > 0:
            # Get the most recent month (assuming sorted_months is chronological)
            recent_month = sorted_months[-1] if isinstance(sorted_months, list) else sorted_months
            # Format as "NUMBER OF BILLED CUSTOMERS - JUL" (just month, no year)
            month_only = f" - {recent_month.upper()}"
        
        # Enhanced title formatting with month only
        formatted_title = f"NUMBER OF BILLED CUSTOMERS{month_only}"
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = formatted_title
        p = title_frame.paragraphs[0]
        p.font.name = "Calibri"  # CHANGED TO CALIBRI
        p.font.size = Pt(28)      # TITLE SIZE = 28PT
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.CENTER
        
        # Use exact DataFrame column order (after S.No removal)
        available_columns = list(df.columns)
        
        num_rows = len(df) + 1
        num_cols = len(available_columns)
        
        # Table positioning and sizing - INCREASED WIDTH
        table_width = Inches(12.5)  # INCREASED from 11.0 to 12.5
        table_height = Inches(len(df) * 0.35 + 0.5)
        table_left = Inches(0.4)    # REDUCED from 1.165 to 0.4
        table_top = Inches(1.3)
        
        table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height).table
        
        # Modified column width distribution (without S.No column) - INCREASED WIDTHS
        if num_cols > 0:
            if num_cols >= 2:
                # First column is now Executive Name, give it more width
                table.columns[0].width = Inches(4.5)  # INCREASED from 4.0 to 4.5 - Executive Name
                remaining_width = 8.0  # INCREASED from 7.0 to 8.0
                if num_cols > 1:
                    month_col_width = remaining_width / (num_cols - 1)
                    for i in range(1, num_cols):
                        table.columns[i].width = Inches(month_col_width)
            else:
                # Only one column case
                col_width = 12.5 / num_cols  # INCREASED from 11.0 to 12.5
                for i in range(num_cols):
                    table.columns[i].width = Inches(col_width)
        
        # Enhanced header row - INCREASED FONT SIZE
        for col_idx, col_name in enumerate(available_columns):
            cell = table.cell(0, col_idx)
            formatted_header = format_column_header(col_name)
            cell.text = formatted_header
            
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Calibri"  # CHANGED TO CALIBRI
            para.font.size = Pt(16)      # INCREASED HEADER SIZE from 15PT to 16PT
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            
            cell.margin_left = Inches(0.03)   # REDUCED margins for more text space
            cell.margin_right = Inches(0.03)  # REDUCED margins for more text space
            cell.margin_top = Inches(0.02)    # REDUCED margins for more text space
            cell.margin_bottom = Inches(0.02) # REDUCED margins for more text space
        
        # FIXED: Enhanced data rows with PERFECT TOTAL ROW STYLING and REDUCED MARGINS
        for row_idx, (index, row) in enumerate(df.iterrows(), start=1):
            # FIXED: Enhanced TOTAL ROW DETECTION
            first_col_value = row.iloc[0]
            is_total = is_total_row(first_col_value)
            
            for col_idx, col_name in enumerate(available_columns):
                cell = table.cell(row_idx, col_idx)
                try:
                    value = row[col_name]
                    cell.text = str(value)
                except (KeyError, ValueError) as e:
                    cell.text = ""
                    logger.warning(f"Error accessing {col_name} at row {index} in customer slide: {e}")
                
                para = cell.text_frame.paragraphs[0]
                para.font.name = "Calibri"  # CHANGED TO CALIBRI
                para.alignment = PP_ALIGN.CENTER
                
                cell.margin_left = Inches(0.03)   # REDUCED margins for more text space
                cell.margin_right = Inches(0.03)  # REDUCED margins for more text space
                cell.margin_top = Inches(0.02)    # REDUCED margins for more text space
                cell.margin_bottom = Inches(0.02) # REDUCED margins for more text space
                
                cell.fill.solid()
                
                # FIXED: PERFECT TOTAL ROW FORMATTING - SAME AS ALL OTHER FUNCTIONS
                if is_total:
                    # Perfect Total Row Styling - Gray background, Calibri, bold white text, same size
                    cell.fill.fore_color.rgb = RGBColor(169, 169, 169)  # Gray background
                    para.font.name = "Calibri"  # Calibri font
                    para.font.bold = True
                    para.font.size = Pt(14)  # DATA SIZE = 14PT
                    para.font.color.rgb = RGBColor(255, 255, 255)  # White text
                else:
                    # Regular data row styling
                    para.font.bold = False
                    para.font.size = Pt(14)  # DATA SIZE = 14PT
                    para.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                    
                    # Alternating row colors
                    if (row_idx - 1) % 2 == 0:
                        cell.fill.fore_color.rgb = RGBColor(221, 235, 247)  # Light blue
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        
    except Exception as e:
        logger.error(f"Error creating customer PPT slide: {e}")
        raise

def create_od_individual_ppt(od_target_df, od_title, logo_file=None, executive_name=None, date_str=None, branch_name=None):
    """Create individual PPT for OD Target report with slide numbers and thank you slide"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Extract month from od_title for simplified formatting
        # Format as "OD TARGET - AUG 2025" (remove the date range part)
        simplified_title = "OD TARGET"
        if od_title and "OD TARGET" in od_title.upper():
            # Extract just the month and year part
            parts = od_title.split('-')
            if len(parts) > 1:
                # Get the month-year part (e.g., "AUG 2025")
                month_year_part = parts[1].strip()
                # Remove any date range in parentheses
                if '(' in month_year_part:
                    month_year_part = month_year_part.split('(')[0].strip()
                simplified_title = f"OD TARGET - {month_year_part}"
        
        # Create title slide with executive information
        create_title_slide(prs, simplified_title, logo_file, executive_name, date_str, branch_name)
        
        # FIXED: Use process_df_for_slides for proper splitting of OD Target data
        if not od_target_df.empty:
            process_df_for_slides(prs, od_target_df, simplified_title, percent_cols=[], is_consolidated=False)
        
        # Finalize with slide numbers and thank you slide
        finalize_presentation_with_numbering_and_thank_you(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        logger.info(f"✅ OD Target PPT created with title: {simplified_title}")
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating OD Target PPT: {e}")
        return None

def create_od_ppt_slide(slide, df, title):
    """Create OD PPT slide with enhanced formatting and FIXED TOTAL ROW STYLING - INCREASED WIDTHS"""
    try:
        # Extract month from title for simplified formatting
        # Format as "OD TARGET - AUG 2025" (remove the date range part)
        simplified_title = "OD TARGET"
        if title and "OD TARGET" in title.upper():
            # Extract just the month and year part
            parts = title.split('-')
            if len(parts) > 1:
                # Get the month-year part (e.g., "AUG 2025")
                month_year_part = parts[1].strip()
                # Remove any date range in parentheses
                if '(' in month_year_part:
                    month_year_part = month_year_part.split('(')[0].strip()
                simplified_title = f"OD TARGET - {month_year_part}"
        
        # Enhanced title formatting
        formatted_title = simplified_title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = formatted_title
        p = title_frame.paragraphs[0]
        p.font.name = "Calibri"  # CHANGED TO CALIBRI
        p.font.size = Pt(28)      # TITLE SIZE = 28PT
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.CENTER
        
        # Use exact DataFrame column order
        available_columns = list(df.columns)
        
        num_rows = len(df) + 1
        num_cols = len(available_columns)
        
        # FIXED TABLE POSITIONING AND SIZING FOR OD TARGET - INCREASED WIDTH
        table_width = Inches(12.5)  # INCREASED from 11.0 to 12.5
        table_height = Inches(len(df) * 0.35 + 0.5)
        table_left = Inches(0.4)    # REDUCED from 1.165 to 0.4
        table_top = Inches(1.3)
        
        table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height).table
        
        # Column width distribution - INCREASED WIDTHS
        if num_cols > 0:
            if num_cols == 2:
                table.columns[0].width = Inches(7.0)   # INCREASED from 6.0 to 7.0
                table.columns[1].width = Inches(5.5)   # INCREASED from 5.0 to 5.5
            else:
                col_width = 12.5 / num_cols  # INCREASED from 11.0 to 12.5
                for i in range(num_cols):
                    table.columns[i].width = Inches(col_width)
        
        # Enhanced header row - INCREASED FONT SIZE AND REDUCED MARGINS
        for i, col_name in enumerate(available_columns):
            header_cell = table.cell(0, i)
            formatted_header = format_column_header(col_name)
            header_cell.text = formatted_header
            
            para = header_cell.text_frame.paragraphs[0]
            para.font.name = "Calibri"  # CHANGED TO CALIBRI
            para.font.bold = True
            para.font.size = Pt(16)      # INCREASED HEADER SIZE from 15PT to 16PT
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            
            header_cell.fill.solid()
            header_cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            
            header_cell.margin_left = Inches(0.03)   # REDUCED margins for more text space
            header_cell.margin_right = Inches(0.03)  # REDUCED margins for more text space
            header_cell.margin_top = Inches(0.02)    # REDUCED margins for more text space
            header_cell.margin_bottom = Inches(0.02) # REDUCED margins for more text space
        
        # FIXED: Enhanced data rows with PERFECT TOTAL ROW STYLING and REDUCED MARGINS
        for i in range(len(df)):
            # FIXED: Enhanced TOTAL ROW DETECTION
            first_col_value = df.iloc[i].iloc[0]
            is_total = is_total_row(first_col_value)
            
            # ENHANCED: Log TOTAL row detection for OD Target debugging
            logger.info(f"🔍 OD Target row {i+1}: '{first_col_value}' -> is_total: {is_total}")
            
            for j, col_name in enumerate(available_columns):
                cell = table.cell(i + 1, j)
                value = df.iloc[i][col_name] if col_name in df.columns else ""
                
                if col_name == 'TARGET':
                    cell.text = f"{float(value):.2f}" if isinstance(value, (int, float)) else str(value)
                else:
                    cell.text = str(value)
                
                para = cell.text_frame.paragraphs[0]
                para.font.name = "Calibri"  # CHANGED TO CALIBRI
                para.alignment = PP_ALIGN.CENTER
                
                cell.margin_left = Inches(0.03)   # REDUCED margins for more text space
                cell.margin_right = Inches(0.03)  # REDUCED margins for more text space
                cell.margin_top = Inches(0.02)    # REDUCED margins for more text space
                cell.margin_bottom = Inches(0.02) # REDUCED margins for more text space
                
                cell.fill.solid()
                
                # ENHANCED: TOTAL ROW FORMATTING with GRAY BACKGROUND for OD Target
                if is_total:
                    # Perfect Total Row Styling - Gray background, Calibri, bold black text
                    cell.fill.fore_color.rgb = RGBColor(169, 169, 169)  # Gray background
                    para.font.name = "Calibri"  # Calibri font
                    para.font.bold = True
                    para.font.size = Pt(14)  # DATA SIZE = 14PT
                    para.font.color.rgb = RGBColor(0, 0, 0)  # Black text for better readability
                    
                    if j == 0:
                        cell.text = cell.text.upper()
                        # ENHANCED: Log when gray background is applied to OD Target TOTAL row
                        logger.info(f"✅ Applied gray background to OD Target TOTAL row: '{first_col_value}'")
                else:
                    # Regular data row styling
                    para.font.bold = False
                    para.font.size = Pt(14)  # DATA SIZE = 14PT
                    para.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                    
                    # Alternating row colors
                    if i % 2 == 0:
                        cell.fill.fore_color.rgb = RGBColor(221, 235, 247)  # Light blue
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        
    except Exception as e:
        logger.error(f"Error creating OD PPT slide: {e}")
        raise

# Utility functions with enhanced formatting support
def validate_ppt_data(results_data):
    """Validate that results data contains the required structure for PPT generation"""
    if not isinstance(results_data, dict):
        raise ValueError("Results data must be a dictionary")
    
    # Check for budget vs billed data
    required_keys = ['budget_vs_billed_qty', 'budget_vs_billed_value', 'overall_sales_qty', 'overall_sales_value']
    has_budget_data = any(key in results_data for key in required_keys)
    
    if has_budget_data:
        valid_data_found = False
        for key in required_keys:
            if key in results_data:
                data = results_data[key]
                if not isinstance(data, list):
                    raise ValueError(f"Data for {key} must be a list")
                if len(data) == 0:
                    logger.warning(f"Data for {key} is empty")
                    continue
                first_item = data[0]
                if not isinstance(first_item, dict) or 'Executive' not in first_item:
                    raise ValueError(f"Data for {key} must contain 'Executive' column")
                valid_data_found = True
        if not valid_data_found:
            raise ValueError("At least one budget vs billed dataset must be non-empty and valid")
    
    # Check for OD data
    elif 'od_results' in results_data:
        od_results = results_data['od_results']
        if not isinstance(od_results, list):
            raise ValueError("od_results must be a list")
        if not od_results:
            raise ValueError("od_results is empty")
        required_columns = [
            'Executive', 'Due Target', 'Collection Achieved', 'Overall % Achieved',
            'For the month Overdue', 'For the month Collection', '% Achieved (Selected Month)'
        ]
        first_row = od_results[0]
        for col in required_columns:
            if col not in first_row:
                raise ValueError(f"Missing required column: {col}")
    
    else:
        raise ValueError("Results data must contain either budget vs billed data or OD results data")
    
    return True

def estimate_slide_count_for_consolidated(dfs_info):
    """Estimate total slide count for consolidated PPT with ALL DATA INCLUDED"""
    try:
        slide_count = 1  # Title slide
        
        for df_info in dfs_info:
            df_data = df_info['df']
            title = df_info['title'].lower()
            
            # Convert to DataFrame if it's a list of dicts
            if isinstance(df_data, list):
                df = pd.DataFrame(df_data)
            else:
                df = df_data.copy()
            
            if df.empty:
                continue
                
            # Count data rows and TOTAL rows separately
            first_col = df.columns[0] if not df.empty else None
            if first_col:
                total_mask = df.iloc[:, 0].apply(is_total_row)
                data_rows_count = len(df[~total_mask]) if total_mask.any() else len(df)
                has_total = total_mask.any()
                
                # Use 14 as split threshold to make room for TOTAL
                split_threshold = 14
                
                if data_rows_count <= split_threshold:
                    slide_count += 1  # Single slide (total will be on same slide)
                else:
                    # Calculate number of parts needed
                    num_parts = math.ceil(data_rows_count / split_threshold)
                    slide_count += num_parts  # All parts (TOTAL on last slide)
            else:
                slide_count += 1
        
        slide_count += 1  # Thank you slide
        return slide_count
        
    except Exception as e:
        logger.error(f"Error estimating slide count: {e}")
        return len(dfs_info) + 2  # Fallback

def get_actual_slide_count(ppt_buffer):
    """Get actual slide count from generated PPT buffer"""
    try:
        ppt_buffer.seek(0)
        temp_prs = Presentation(ppt_buffer)
        slide_count = len(temp_prs.slides)
        ppt_buffer.seek(0)  # Reset buffer position
        return slide_count
    except Exception as e:
        logger.error(f"Error getting slide count: {e}")
        return 0

def log_ppt_generation_summary(report_type, slide_count, report_count=None):
    """Log PPT generation summary"""
    if report_count:
        logger.info(f"📊 {report_type} PPT generated: {report_count} reports → {slide_count} slides (ALL DATA INCLUDED)")
        print(f"📊 {report_type} PPT generated: {report_count} reports → {slide_count} slides (ALL DATA INCLUDED)")
    else:
        logger.info(f"📊 {report_type} PPT generated with {slide_count} slides (ALL DATA INCLUDED)")
        print(f"📊 {report_type} PPT generated with {slide_count} slides (ALL DATA INCLUDED)")

def clean_dataframe_for_ppt(df):
    """Clean DataFrame before PPT generation with enhanced formatting"""
    try:
        # Remove any completely empty rows
        df = df.dropna(how='all')
        
        # Fill NaN values with empty strings for text columns
        for col in df.columns:
            if df[col].dtype == 'object':
                df[col] = df[col].fillna('')
            else:
                df[col] = df[col].fillna(0)
        
        # Round numeric columns
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            df[col] = df[col].round(2)
        
        return df
        
    except Exception as e:
        logger.error(f"Error cleaning DataFrame: {e}")
        return df

def validate_product_growth_ppt_data(results_data):
    """Validate data structure for Product Growth PPT generation"""
    if not isinstance(results_data, dict):
        raise ValueError("Results data must be a dictionary")
    
    if not results_data:
        raise ValueError("Results data is empty")
    
    # Check if each company has required data structure
    required_keys = ['qty_df', 'value_df']
    required_columns_qty = ['PRODUCT GROUP', 'LY_QTY', 'BUDGET_QTY', 'CY_QTY', 'ACHIEVEMENT %']
    required_columns_value = ['PRODUCT GROUP', 'LY_VALUE', 'BUDGET_VALUE', 'CY_VALUE', 'ACHIEVEMENT %']
    
    for company, data in results_data.items():
        for key in required_keys:
            if key not in data:
                raise ValueError(f"Missing required key '{key}' for company '{company}'")
            
            df_data = data[key]
            if not isinstance(df_data, list) or not df_data:
                raise ValueError(f"Data for '{key}' in company '{company}' must be a non-empty list")
            
            # Check columns
            first_row = df_data[0]
            if key == 'qty_df':
                required_cols = required_columns_qty
            else:
                required_cols = required_columns_value
            
            for col in required_cols:
                if col not in first_row:
                    raise ValueError(f"Missing required column '{col}' in {key} for company '{company}'")
    
    return True


# Export all enhanced functions
__all__ = [
    'add_slide_number',
    'create_thank_you_slide', 
    'finalize_presentation_with_numbering_and_thank_you',
    'format_column_header',
    'format_slide_title',
    'create_title_slide',
    'add_table_slide', 
    'create_product_growth_ppt',
    'process_df_for_slides',
    'estimate_slide_count_for_consolidated',
    'create_consolidated_ppt',
    'create_customer_ppt_slide',
    'create_nbc_individual_ppt',
    'create_od_ppt_slide',
    'create_od_individual_ppt',
    'create_executive_budget_ppt',
    'create_executive_od_ppt',
    'fix_budget_vs_billed_column_order',
    'create_budget_vs_billed_ppt_with_fixed_order',
    'validate_ppt_data',
    'validate_product_growth_ppt_data',
    'get_actual_slide_count',
    'log_ppt_generation_summary',
    'clean_dataframe_for_ppt',
    'calculate_slide_constraints',
    'calculate_max_fitting_rows',
    'is_total_row'
]